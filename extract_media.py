"""
Stage 1: Media Extraction
Extracts images from PPTX files using ZIP/XML parsing.
Maps each image to its slide number.
(Excludes ffmpeg/video processing — not available on GitHub Actions.)
"""

import sys
import os
import json
import re
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from datetime import datetime, timezone

sys.stdout.reconfigure(encoding="utf-8")

from config import SOURCES_DIR, MEDIA_DIR, LOGS_DIR

# XML namespaces in PPTX
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".emf", ".wmf", ".svg"}


def build_slide_image_mapping(zip_ref):
    """Parse PPTX XML to map media filenames → slide numbers."""
    media_to_slides = {}
    slide_files = sorted(
        f for f in zip_ref.namelist()
        if re.match(r"ppt/slides/slide\d+\.xml$", f)
    )

    for slide_file in slide_files:
        slide_num = int(re.search(r"slide(\d+)\.xml", slide_file).group(1))
        rels_file = slide_file.replace("slides/", "slides/_rels/") + ".rels"

        if rels_file not in zip_ref.namelist():
            continue

        # Parse rels to get rId → media filename
        rels_root = ET.fromstring(zip_ref.read(rels_file).decode("utf-8"))
        rid_to_media = {}
        for rel in rels_root.findall(
            ".//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship"
        ):
            target = rel.get("Target", "")
            rel_type = rel.get("Type", "")
            if "image" in rel_type.lower() or "../media/" in target:
                rid_to_media[rel.get("Id")] = Path(target).name

        # Parse slide XML to find which rIds are referenced
        slide_root = ET.fromstring(zip_ref.read(slide_file).decode("utf-8"))
        for elem in slide_root.iter():
            for attr in ("embed", "link"):
                rid = elem.get(f"{{{NS['r']}}}{attr}")
                if rid and rid in rid_to_media:
                    media_name = rid_to_media[rid]
                    media_to_slides.setdefault(media_name, [])
                    if slide_num not in media_to_slides[media_name]:
                        media_to_slides[media_name].append(slide_num)

    for v in media_to_slides.values():
        v.sort()

    return media_to_slides


def extract_pptx_images(pptx_path, output_base):
    """Extract all images from a single PPTX file."""
    safe_name = pptx_path.stem.replace(" ", "_")
    output_folder = output_base / safe_name
    output_folder.mkdir(parents=True, exist_ok=True)

    extracted = []

    try:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            media_to_slides = build_slide_image_mapping(zf)
            media_files = [f for f in zf.namelist() if f.startswith("ppt/media/")]

            for i, media_file in enumerate(media_files):
                ext = Path(media_file).suffix.lower()
                if ext not in IMAGE_EXTS:
                    continue

                data = zf.read(media_file)
                media_name = Path(media_file).name
                slide_numbers = media_to_slides.get(media_name, [])

                output_name = f"image_{i + 1:03d}{ext}"
                output_path = output_folder / output_name

                with open(output_path, "wb") as f:
                    f.write(data)

                extracted.append({
                    "source_pptx": str(pptx_path),
                    "image_path": str(output_path),
                    "original_media_name": media_name,
                    "index": i + 1,
                    "slide_numbers": slide_numbers,
                    "primary_slide": slide_numbers[0] if slide_numbers else None,
                    "size_bytes": len(data),
                    "extension": ext,
                })

    except Exception as e:
        print(f"  Could not extract media from {pptx_path.name}: {e}"
              f" — text content still available via native API extraction")

    return extracted


def _iter_shapes_recursive(shapes):
    """Recursively yield all shapes, descending into GroupShapes."""
    from pptx.shapes.group import GroupShape
    for shape in shapes:
        yield shape
        if isinstance(shape, GroupShape):
            yield from _iter_shapes_recursive(shape.shapes)


def extract_pptx_links(pptx_path):
    """Extract hyperlinks from a PPTX file using python-pptx.
    Recursively descends into GroupShapes.
    Returns list of {url, text, slide_number, link_type}."""
    links = []
    try:
        from pptx import Presentation

        prs = Presentation(pptx_path)

        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in _iter_shapes_recursive(slide.shapes):

                # Text hyperlinks from text frames
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                if run.hyperlink and run.hyperlink.address:
                                    links.append({
                                        "url": run.hyperlink.address,
                                        "text": run.text.strip(),
                                        "slide_number": slide_num,
                                        "link_type": "text_hyperlink",
                                    })
                            except (KeyError, Exception):
                                pass

                # Click action hyperlinks on shapes
                try:
                    if hasattr(shape, "click_action") and shape.click_action and shape.click_action.hyperlink:
                        addr = shape.click_action.hyperlink.address
                        if addr and (addr.startswith("http") or addr.startswith("mailto:")):
                            links.append({
                                "url": addr,
                                "text": shape.text.strip() if shape.has_text_frame else "",
                                "slide_number": slide_num,
                                "link_type": "click_action",
                            })
                except Exception:
                    pass

    except Exception as e:
        print(f"  Could not extract links from {pptx_path.name}: {e} — links may be available from native API")

    return links


def extract_docx_links(docx_path):
    """Extract hyperlinks from a DOCX file by parsing relationship XML.
    Returns list of {url, text, link_type}."""
    links = []
    try:
        from docx import Document

        doc = Document(str(docx_path))

        # Build rId → URL mapping from relationships
        rels = doc.part.rels
        rid_to_url = {}
        for rel_id, rel in rels.items():
            if "hyperlink" in rel.reltype:
                rid_to_url[rel_id] = rel.target_ref

        # Parse document XML to find hyperlink elements
        nsmap = {
            "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
            "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        }
        body = doc.element.body
        for para_elem in body.findall(".//w:p", nsmap):
            for hlink in para_elem.findall(".//w:hyperlink", nsmap):
                rid = hlink.get(f"{{{nsmap['r']}}}id", "")
                url = rid_to_url.get(rid, "")
                if not url:
                    continue
                # Skip local file paths — only keep web URLs
                if not url.startswith(("http://", "https://", "mailto:")):
                    continue

                # Get the visible text within the hyperlink
                texts = []
                for run_elem in hlink.findall(".//w:t", nsmap):
                    if run_elem.text:
                        texts.append(run_elem.text)
                link_text = "".join(texts).strip()

                links.append({
                    "url": url,
                    "text": link_text,
                    "link_type": "docx_hyperlink",
                })

    except Exception as e:
        print(f"  Could not extract links from {docx_path.name}: {e} — non-critical")

    return links


def extract_pdf_links(pdf_path):
    """Extract annotation hyperlinks from a PDF file.
    Returns list of {url, page_number}."""
    links = []
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(pdf_path))

        for page_num, page in enumerate(reader.pages, 1):
            annots = page.get("/Annots")
            if not annots:
                continue
            annot_list = annots.get_object() if hasattr(annots, "get_object") else annots
            if not isinstance(annot_list, (list, tuple)):
                continue

            for annot_ref in annot_list:
                try:
                    annot = annot_ref.get_object() if hasattr(annot_ref, "get_object") else annot_ref
                    if not isinstance(annot, dict):
                        continue
                    if str(annot.get("/Subtype", "")) != "/Link":
                        continue
                    action = annot.get("/A")
                    if action:
                        action_obj = action.get_object() if hasattr(action, "get_object") else action
                        if isinstance(action_obj, dict):
                            uri = str(action_obj.get("/URI", ""))
                            if uri and uri.startswith("http"):
                                links.append({
                                    "url": uri,
                                    "page_number": page_num,
                                    "link_type": "pdf_annotation",
                                })
                except Exception:
                    continue

    except Exception as e:
        print(f"  Could not extract links from {pdf_path.name}: {e} — non-critical")

    return links


def run_extraction(source_dir=None):
    """Run media extraction on all source files."""
    print("=" * 60)
    print("  Stage 1: Media Extraction")
    print("=" * 60)
    print()

    src = Path(source_dir) if source_dir else SOURCES_DIR
    MEDIA_DIR.mkdir(parents=True, exist_ok=True)

    pptx_files = list(src.rglob("*.pptx"))
    # Skip temp files
    pptx_files = [f for f in pptx_files if not f.name.startswith("~$")]
    print(f"Found {len(pptx_files)} PPTX files\n")

    all_results = {
        "extracted_at": datetime.now(timezone.utc).isoformat(),
        "pptx_files": [],
        "total_images": 0,
        "total_with_slides": 0,
        "total_without_slides": 0,
        "total_links": 0,
    }

    for pptx_path in pptx_files:
        print(f"Processing: {pptx_path.name}")
        images = extract_pptx_images(pptx_path, MEDIA_DIR)
        links = extract_pptx_links(pptx_path)

        with_slides = sum(1 for img in images if img["slide_numbers"])
        without_slides = len(images) - with_slides

        print(f"  Extracted {len(images)} images ({with_slides} with slide mapping), {len(links)} links")

        all_results["pptx_files"].append({
            "source": str(pptx_path),
            "relative_path": str(pptx_path.relative_to(src)) if src in pptx_path.parents else pptx_path.name,
            "images_count": len(images),
            "with_slide_mapping": with_slides,
            "images": images,
            "links_count": len(links),
            "links": links,
        })
        all_results["total_images"] += len(images)
        all_results["total_with_slides"] += with_slides
        all_results["total_without_slides"] += without_slides
        all_results["total_links"] += len(links)

    # ── DOCX link extraction ──
    docx_files = list(src.rglob("*.docx"))
    docx_files = [f for f in docx_files if not f.name.startswith("~$")]
    print(f"\nFound {len(docx_files)} DOCX files")

    all_results["docx_files"] = []
    all_results["total_docx_links"] = 0

    for docx_path in docx_files:
        links = extract_docx_links(docx_path)
        if links:
            try:
                rel = str(docx_path.relative_to(src))
            except ValueError:
                rel = docx_path.name
            print(f"  {docx_path.name}: {len(links)} links")
            all_results["docx_files"].append({
                "source": str(docx_path),
                "relative_path": rel,
                "links_count": len(links),
                "links": links,
            })
            all_results["total_docx_links"] += len(links)
            all_results["total_links"] += len(links)

    print(f"Total DOCX links extracted: {all_results['total_docx_links']}")

    # ── PDF link extraction ──
    pdf_files = list(src.rglob("*.pdf"))
    print(f"\nFound {len(pdf_files)} PDF files")

    all_results["pdf_files"] = []
    all_results["total_pdf_links"] = 0

    for pdf_path in pdf_files:
        links = extract_pdf_links(pdf_path)
        if links:
            try:
                rel = str(pdf_path.relative_to(src))
            except ValueError:
                rel = pdf_path.name
            print(f"  {pdf_path.name}: {len(links)} links")
            all_results["pdf_files"].append({
                "source": str(pdf_path),
                "relative_path": rel,
                "links_count": len(links),
                "links": links,
            })
            all_results["total_pdf_links"] += len(links)
            all_results["total_links"] += len(links)

    print(f"Total PDF links extracted: {all_results['total_pdf_links']}")

    # Save metadata
    metadata_path = MEDIA_DIR / "extraction_metadata.json"
    with open(metadata_path, "w", encoding="utf-8") as f:
        json.dump(all_results, f, indent=2, ensure_ascii=False)

    print(f"\nTotal images extracted: {all_results['total_images']}")
    print(f"  With slide mapping: {all_results['total_with_slides']}")
    print(f"  Without: {all_results['total_without_slides']}")
    print(f"Total links extracted: {all_results['total_links']}")
    print(f"  PPTX: {all_results['total_links'] - all_results['total_docx_links'] - all_results['total_pdf_links']}")
    print(f"  DOCX: {all_results['total_docx_links']}")
    print(f"  PDF:  {all_results['total_pdf_links']}")
    print(f"Metadata saved: {metadata_path}")
    print("=" * 60)

    return all_results


if __name__ == "__main__":
    run_extraction()
