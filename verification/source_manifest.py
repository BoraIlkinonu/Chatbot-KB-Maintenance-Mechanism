"""
Source Manifest Builder.

Parses PPTX, DOCX, PDF, and native Google extracts directly (bypassing the pipeline)
to build a ground-truth list of every ContentAtom in the source files.

Key improvement over pipeline: recursively descends into GroupShapes and extracts
hyperlinks from speaker notes.
"""

import json
import re
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

from verification import ContentAtom, SourceManifest

# Reuse path-parsing functions from consolidate
from consolidate import extract_term_from_path, extract_lesson_from_path

from config import BASE_DIR

# Categories that count toward coverage metrics
_LESSON_CATEGORIES = {"lesson_content", "pending_review"}
# Categories excluded from coverage
_EXCLUDED_CATEGORIES = {"admin_doc", "support_resource", "media"}

# XML namespaces used in PPTX
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".emf", ".wmf", ".svg"}

# Max lessons before collapsing to lesson=None to avoid atom duplication explosion.
# Files spanning >3 lessons (e.g., assessment guides covering all 24 lessons) would
# create N copies of every atom. Instead, store once with lesson=None.
_MAX_LESSON_DUP = 3


def _lesson_iter(lessons):
    """Return lesson list for atom creation. Collapse large lists to [None]."""
    if not lessons:
        return [None]
    if len(lessons) <= _MAX_LESSON_DUP:
        return lessons
    return [None]


# ──────────────────────────────────────────────────────────
# PPTX parsing (with GroupShape recursion)
# ──────────────────────────────────────────────────────────

def _iter_shapes_recursive(shapes):
    """Recursively yield all shapes, descending into GroupShapes."""
    from pptx.shapes.group import GroupShape
    for shape in shapes:
        yield shape
        if isinstance(shape, GroupShape):
            yield from _iter_shapes_recursive(shape.shapes)


def _extract_pptx_atoms(pptx_path, rel_path, term, lessons):
    """Extract all ContentAtoms from a PPTX file."""
    atoms = []
    try:
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
    except Exception as e:
        print(f"  Warning: cannot open PPTX {pptx_path.name}: {e}")
        return atoms

    for slide_num, slide in enumerate(prs.slides, 1):
        location_prefix = f"slide:{slide_num}"

        # Iterate all shapes recursively (fixes GroupShape skip)
        for shape in _iter_shapes_recursive(slide.shapes):

            # Text content from text frames
            if shape.has_text_frame:
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    para_text = para.text.strip()
                    if para_text:
                        for lesson in _lesson_iter(lessons):
                            atoms.append(ContentAtom(
                                atom_type="text_block",
                                content=para_text,
                                source_file=rel_path,
                                location=f"{location_prefix}:para:{para_idx}",
                                term=term,
                                lesson=lesson,
                                metadata={"shape_name": getattr(shape, "name", "")},
                            ))

                    # Hyperlinks from runs
                    for run in para.runs:
                        try:
                            if run.hyperlink and run.hyperlink.address:
                                url = run.hyperlink.address
                                for lesson in _lesson_iter(lessons):
                                    atoms.append(ContentAtom(
                                        atom_type="link",
                                        content=url,
                                        source_file=rel_path,
                                        location=f"{location_prefix}:run_hyperlink",
                                        term=term,
                                        lesson=lesson,
                                        metadata={"link_text": run.text.strip()},
                                    ))
                        except (KeyError, Exception):
                            pass

            # Click action hyperlinks
            try:
                if hasattr(shape, "click_action") and shape.click_action and shape.click_action.hyperlink:
                    addr = shape.click_action.hyperlink.address
                    if addr:
                        for lesson in _lesson_iter(lessons):
                            atoms.append(ContentAtom(
                                atom_type="link",
                                content=addr,
                                source_file=rel_path,
                                location=f"{location_prefix}:click_action",
                                term=term,
                                lesson=lesson,
                                metadata={"link_text": shape.text.strip() if shape.has_text_frame else ""},
                            ))
            except Exception:
                pass

            # Tables
            if shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    for cell_idx, cell in enumerate(row.cells):
                        cell_text = cell.text.strip()
                        if cell_text:
                            for lesson in _lesson_iter(lessons):
                                atoms.append(ContentAtom(
                                    atom_type="table",
                                    content=cell_text,
                                    source_file=rel_path,
                                    location=f"{location_prefix}:table:r{row_idx}c{cell_idx}",
                                    term=term,
                                    lesson=lesson,
                                ))

        # Speaker notes — text AND hyperlinks
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = notes_tf.text.strip()
            if notes_text:
                for lesson in _lesson_iter(lessons):
                    atoms.append(ContentAtom(
                        atom_type="speaker_note",
                        content=notes_text,
                        source_file=rel_path,
                        location=f"{location_prefix}:notes",
                        term=term,
                        lesson=lesson,
                    ))

            # Hyperlinks within speaker notes
            for para in notes_tf.paragraphs:
                for run in para.runs:
                    try:
                        if run.hyperlink and run.hyperlink.address:
                            for lesson in _lesson_iter(lessons):
                                atoms.append(ContentAtom(
                                    atom_type="link",
                                    content=run.hyperlink.address,
                                    source_file=rel_path,
                                    location=f"{location_prefix}:notes_hyperlink",
                                    term=term,
                                    lesson=lesson,
                                    metadata={"link_text": run.text.strip(), "in_notes": True},
                                ))
                    except (KeyError, Exception):
                        pass

    # Images via ZIP/XML mapping (reuse approach from extract_media.py)
    atoms.extend(_extract_pptx_images(pptx_path, rel_path, term, lessons))

    return atoms


def _extract_pptx_images(pptx_path, rel_path, term, lessons):
    """Extract image atoms from PPTX using ZIP parsing."""
    atoms = []
    try:
        with zipfile.ZipFile(str(pptx_path), "r") as zf:
            # Build media-to-slide mapping
            media_to_slides = {}
            slide_files = sorted(
                f for f in zf.namelist()
                if re.match(r"ppt/slides/slide\d+\.xml$", f)
            )
            for slide_file in slide_files:
                slide_num = int(re.search(r"slide(\d+)\.xml", slide_file).group(1))
                rels_file = slide_file.replace("slides/", "slides/_rels/") + ".rels"
                if rels_file not in zf.namelist():
                    continue
                rels_root = ET.fromstring(zf.read(rels_file).decode("utf-8"))
                for rel in rels_root.findall(
                    ".//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship"
                ):
                    target = rel.get("Target", "")
                    rel_type = rel.get("Type", "")
                    if "image" in rel_type.lower() or "../media/" in target:
                        media_name = Path(target).name
                        media_to_slides.setdefault(media_name, [])
                        if slide_num not in media_to_slides[media_name]:
                            media_to_slides[media_name].append(slide_num)

            # Count images
            media_files = [f for f in zf.namelist() if f.startswith("ppt/media/")]
            for media_file in media_files:
                ext = Path(media_file).suffix.lower()
                if ext not in IMAGE_EXTS:
                    continue
                media_name = Path(media_file).name
                slide_nums = media_to_slides.get(media_name, [])
                location = f"slide:{slide_nums[0]}" if slide_nums else "unknown_slide"
                for lesson in _lesson_iter(lessons):
                    atoms.append(ContentAtom(
                        atom_type="image",
                        content=media_name,
                        source_file=rel_path,
                        location=f"{location}:image:{media_name}",
                        term=term,
                        lesson=lesson,
                        metadata={"slide_numbers": sorted(slide_nums), "media_name": media_name},
                    ))
    except Exception as e:
        print(f"  Warning: cannot extract images from {pptx_path.name}: {e}")

    return atoms


# ──────────────────────────────────────────────────────────
# DOCX parsing (with run-level hyperlinks)
# ──────────────────────────────────────────────────────────

def _extract_docx_atoms(docx_path, rel_path, term, lessons):
    """Extract all ContentAtoms from a DOCX file."""
    atoms = []
    try:
        from docx import Document
        doc = Document(str(docx_path))
    except Exception as e:
        print(f"  Warning: cannot open DOCX {docx_path.name}: {e}")
        return atoms

    for para_idx, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if text:
            for lesson in _lesson_iter(lessons):
                atoms.append(ContentAtom(
                    atom_type="text_block",
                    content=text,
                    source_file=rel_path,
                    location=f"paragraph:{para_idx}",
                    term=term,
                    lesson=lesson,
                ))

        # Run-level hyperlinks (the pipeline loses these)
        for run in para.runs:
            # python-docx doesn't expose hyperlinks directly on runs in older versions,
            # but we can check the XML for hyperlink elements
            pass

    # Check paragraph XML for hyperlinks
    atoms.extend(_extract_docx_hyperlinks(doc, rel_path, term, lessons))

    # Tables
    for table_idx, table in enumerate(doc.tables):
        for row_idx, row in enumerate(table.rows):
            for cell_idx, cell in enumerate(row.cells):
                cell_text = cell.text.strip()
                if cell_text:
                    for lesson in _lesson_iter(lessons):
                        atoms.append(ContentAtom(
                            atom_type="table",
                            content=cell_text,
                            source_file=rel_path,
                            location=f"table:{table_idx}:r{row_idx}c{cell_idx}",
                            term=term,
                            lesson=lesson,
                        ))

    return atoms


def _extract_docx_hyperlinks(doc, rel_path, term, lessons):
    """Extract hyperlinks from DOCX by parsing the underlying XML."""
    atoms = []
    nsmap = {
        "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }

    try:
        # Get the relationship mapping for the main document part
        rels = doc.part.rels
        rid_to_url = {}
        for rel_id, rel in rels.items():
            if "hyperlink" in rel.reltype:
                rid_to_url[rel_id] = rel.target_ref

        # Find hyperlink elements in the document XML
        body = doc.element.body
        for para_idx, para_elem in enumerate(body.findall(".//w:p", nsmap)):
            for hlink in para_elem.findall(".//w:hyperlink", nsmap):
                rid = hlink.get(f"{{{nsmap['r']}}}id", "")
                url = rid_to_url.get(rid, "")
                # Get the text within the hyperlink
                texts = []
                for run_elem in hlink.findall(".//w:t", nsmap):
                    if run_elem.text:
                        texts.append(run_elem.text)
                link_text = "".join(texts).strip()

                if url:
                    for lesson in _lesson_iter(lessons):
                        atoms.append(ContentAtom(
                            atom_type="link",
                            content=url,
                            source_file=rel_path,
                            location=f"paragraph:{para_idx}:hyperlink",
                            term=term,
                            lesson=lesson,
                            metadata={"link_text": link_text},
                        ))
    except Exception as e:
        print(f"  Warning: DOCX hyperlink extraction failed for {rel_path}: {e}")

    return atoms


# ──────────────────────────────────────────────────────────
# PDF parsing
# ──────────────────────────────────────────────────────────

def _extract_pdf_atoms(pdf_path, rel_path, term, lessons):
    """Extract ContentAtoms from a PDF file."""
    atoms = []
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(pdf_path))
    except Exception as e:
        print(f"  Warning: cannot open PDF {pdf_path.name}: {e}")
        return atoms

    for page_num, page in enumerate(reader.pages, 1):
        # Page text
        text = page.extract_text()
        if text and text.strip():
            for lesson in _lesson_iter(lessons):
                atoms.append(ContentAtom(
                    atom_type="text_block",
                    content=text.strip(),
                    source_file=rel_path,
                    location=f"page:{page_num}",
                    term=term,
                    lesson=lesson,
                ))

        # Annotation links
        annots = page.get("/Annots")
        if annots:
            annot_list = annots.get_object() if hasattr(annots, "get_object") else annots
            if isinstance(annot_list, (list, tuple)):
                for annot_ref in annot_list:
                    try:
                        annot = annot_ref.get_object() if hasattr(annot_ref, "get_object") else annot_ref
                        if not isinstance(annot, dict):
                            continue
                        if str(annot.get("/Subtype", "")) != "/Link":
                            continue
                        action = annot.get("/A")
                        if action:
                            action_obj = action.get_object() if hasattr(action, "get_object") else action
                            if isinstance(action_obj, dict):
                                uri = str(action_obj.get("/URI", ""))
                                if uri:
                                    for lesson in _lesson_iter(lessons):
                                        atoms.append(ContentAtom(
                                            atom_type="link",
                                            content=uri,
                                            source_file=rel_path,
                                            location=f"page:{page_num}:annotation",
                                            term=term,
                                            lesson=lesson,
                                        ))
                    except Exception:
                        continue

        # XObject images (count only)
        try:
            resources = page.get("/Resources")
            if resources:
                res_obj = resources.get_object() if hasattr(resources, "get_object") else resources
                if isinstance(res_obj, dict):
                    xobjects = res_obj.get("/XObject")
                    if xobjects:
                        xobj_dict = xobjects.get_object() if hasattr(xobjects, "get_object") else xobjects
                        if isinstance(xobj_dict, dict):
                            for obj_name in xobj_dict:
                                try:
                                    xobj = xobj_dict[obj_name]
                                    obj = xobj.get_object() if hasattr(xobj, "get_object") else xobj
                                    if hasattr(obj, "get") and str(obj.get("/Subtype", "")) == "/Image":
                                        for lesson in _lesson_iter(lessons):
                                            atoms.append(ContentAtom(
                                                atom_type="image",
                                                content=str(obj_name),
                                                source_file=rel_path,
                                                location=f"page:{page_num}:xobject:{obj_name}",
                                                term=term,
                                                lesson=lesson,
                                                metadata={"source": "pdf"},
                                            ))
                                except Exception:
                                    continue
        except Exception:
            pass

    return atoms


# ──────────────────────────────────────────────────────────
# Native Google extract parsing
# ──────────────────────────────────────────────────────────

def _extract_native_atoms(json_path, rel_path, term, lessons):
    """Extract ContentAtoms from a native Google API JSON extract.

    Handles two formats:
    1. Combined file with 'extractions' array (native_extractions.json)
    2. Individual extraction files with 'type' field
    """
    atoms = []
    try:
        with open(json_path, "r", encoding="utf-8") as f:
            data = json.load(f)
    except Exception as e:
        print(f"  Warning: cannot read native extract {json_path.name}: {e}")
        return atoms

    # Combined extractions file — iterate each extraction with its own term/lesson
    if "extractions" in data:
        for extraction in data["extractions"]:
            source_path = extraction.get("source_path", "")
            ext_term = extract_term_from_path(source_path) if source_path else term
            ext_lessons = extract_lesson_from_path(source_path, ext_term) if source_path else lessons
            if not ext_lessons:
                ext_lessons = [None]
            native_type = extraction.get("native_type", "")
            ext_rel = f"native:{source_path}" if source_path else rel_path
            atoms.extend(_extract_single_native(extraction, ext_rel, ext_term, ext_lessons, native_type))
        return atoms

    # Individual extraction file
    extract_type = data.get("type", "") or data.get("native_type", "")
    atoms.extend(_extract_single_native(data, rel_path, term, lessons, extract_type))
    return atoms


def _extract_single_native(data, rel_path, term, lessons, native_type):
    """Extract atoms from a single native extraction entry.

    Handles actual data formats from extract_native_google.py:
    - Google Docs: content_blocks[{text, style}], links[{url, text}]
    - Google Slides: slides[{texts:[], links:[], videos:[], tables:[{headers,rows}],
                            image_urls:[], speaker_notes:""}]
    - Google Sheets: sheets[{sheet_name, headers, rows:[[]], total_rows}]
    """
    atoms = []

    # Google Docs
    if native_type == "google_doc":
        for idx, block in enumerate(data.get("content_blocks", [])):
            text = block.get("text", "").strip()
            if text:
                for lesson in _lesson_iter(lessons):
                    atoms.append(ContentAtom(
                        atom_type="text_block",
                        content=text,
                        source_file=rel_path,
                        location=f"block:{idx}",
                        term=term,
                        lesson=lesson,
                        metadata={"style": block.get("style", "")},
                    ))

        for idx, link in enumerate(data.get("links", [])):
            url = link.get("url", "")
            if url:
                for lesson in _lesson_iter(lessons):
                    atoms.append(ContentAtom(
                        atom_type="link",
                        content=url,
                        source_file=rel_path,
                        location=f"link:{idx}",
                        term=term,
                        lesson=lesson,
                        metadata={"link_text": link.get("text", "")},
                    ))

    # Google Slides — actual format uses "texts" (list of strings), not "elements"
    elif native_type == "google_slides":
        for slide_idx, slide in enumerate(data.get("slides", [])):
            slide_num = slide.get("slide_number", slide_idx + 1)

            # Text content: "texts" is a flat list of strings
            for text_idx, text in enumerate(slide.get("texts", [])):
                text = str(text).strip()
                if text:
                    for lesson in _lesson_iter(lessons):
                        atoms.append(ContentAtom(
                            atom_type="text_block",
                            content=text,
                            source_file=rel_path,
                            location=f"slide:{slide_num}:text:{text_idx}",
                            term=term,
                            lesson=lesson,
                        ))

            # Fallback: also check "elements" in case of alternate format
            for elem_idx, elem in enumerate(slide.get("elements", [])):
                text = elem.get("text", "").strip() if isinstance(elem, dict) else str(elem).strip()
                if text:
                    for lesson in _lesson_iter(lessons):
                        atoms.append(ContentAtom(
                            atom_type="text_block",
                            content=text,
                            source_file=rel_path,
                            location=f"slide:{slide_num}:elem:{elem_idx}",
                            term=term,
                            lesson=lesson,
                        ))

            # Per-slide links
            for link_idx, link in enumerate(slide.get("links", [])):
                url = link.get("url", "") if isinstance(link, dict) else str(link)
                link_text = link.get("text", "") if isinstance(link, dict) else ""
                if url:
                    for lesson in _lesson_iter(lessons):
                        atoms.append(ContentAtom(
                            atom_type="link",
                            content=url,
                            source_file=rel_path,
                            location=f"slide:{slide_num}:link:{link_idx}",
                            term=term,
                            lesson=lesson,
                            metadata={"link_text": link_text},
                        ))

            # Per-slide videos
            for vid_idx, video in enumerate(slide.get("videos", [])):
                url = video.get("url", "") if isinstance(video, dict) else str(video)
                if url:
                    for lesson in _lesson_iter(lessons):
                        atoms.append(ContentAtom(
                            atom_type="video_url",
                            content=url,
                            source_file=rel_path,
                            location=f"slide:{slide_num}:video:{vid_idx}",
                            term=term,
                            lesson=lesson,
                        ))

            # Per-slide tables: {headers:[], rows:[[]]}
            for table_idx, table in enumerate(slide.get("tables", [])):
                if isinstance(table, dict):
                    for h_idx, h in enumerate(table.get("headers", [])):
                        if str(h).strip():
                            for lesson in _lesson_iter(lessons):
                                atoms.append(ContentAtom(
                                    atom_type="table",
                                    content=str(h).strip(),
                                    source_file=rel_path,
                                    location=f"slide:{slide_num}:table:{table_idx}:h{h_idx}",
                                    term=term,
                                    lesson=lesson,
                                ))
                    for row_idx, row in enumerate(table.get("rows", [])):
                        for col_idx, cell in enumerate(row):
                            if str(cell).strip():
                                for lesson in _lesson_iter(lessons):
                                    atoms.append(ContentAtom(
                                        atom_type="table",
                                        content=str(cell).strip(),
                                        source_file=rel_path,
                                        location=f"slide:{slide_num}:table:{table_idx}:r{row_idx}c{col_idx}",
                                        term=term,
                                        lesson=lesson,
                                    ))

            # Per-slide image URLs — separate type from extracted PPTX images
            for img_idx, img_url in enumerate(slide.get("image_urls", [])):
                if str(img_url).strip():
                    for lesson in _lesson_iter(lessons):
                        atoms.append(ContentAtom(
                            atom_type="image_ref",
                            content=str(img_url).strip(),
                            source_file=rel_path,
                            location=f"slide:{slide_num}:image_url:{img_idx}",
                            term=term,
                            lesson=lesson,
                            metadata={"source": "native_slides"},
                        ))

            # Speaker notes
            notes = slide.get("speaker_notes", "").strip()
            if notes:
                for lesson in _lesson_iter(lessons):
                    atoms.append(ContentAtom(
                        atom_type="speaker_note",
                        content=notes,
                        source_file=rel_path,
                        location=f"slide:{slide_num}:notes",
                        term=term,
                        lesson=lesson,
                    ))

        # Top-level links/videos (fallback — some formats may have these)
        for idx, video in enumerate(data.get("videos", [])):
            url = video.get("url", "") if isinstance(video, dict) else str(video)
            if url:
                for lesson in _lesson_iter(lessons):
                    atoms.append(ContentAtom(
                        atom_type="video_url",
                        content=url,
                        source_file=rel_path,
                        location=f"video:{idx}",
                        term=term,
                        lesson=lesson,
                    ))

        for idx, link in enumerate(data.get("links", [])):
            url = link.get("url", "") if isinstance(link, dict) else str(link)
            if url:
                for lesson in _lesson_iter(lessons):
                    atoms.append(ContentAtom(
                        atom_type="link",
                        content=url,
                        source_file=rel_path,
                        location=f"link:{idx}",
                        term=term,
                        lesson=lesson,
                    ))

    # Google Sheets — actual format uses "rows"/"sheet_name", not "data"/"name"
    elif native_type == "google_sheet":
        for sheet in data.get("sheets", []):
            sheet_name = sheet.get("sheet_name", "") or sheet.get("name", "")

            # Headers
            for h_idx, h in enumerate(sheet.get("headers", [])):
                if str(h).strip():
                    for lesson in _lesson_iter(lessons):
                        atoms.append(ContentAtom(
                            atom_type="table",
                            content=str(h).strip(),
                            source_file=rel_path,
                            location=f"sheet:{sheet_name}:h{h_idx}",
                            term=term,
                            lesson=lesson,
                        ))

            # Rows (actual key) + data (fallback key)
            rows = sheet.get("rows", []) or sheet.get("data", [])
            for row_idx, row in enumerate(rows):
                if not isinstance(row, (list, tuple)):
                    continue
                for col_idx, cell in enumerate(row):
                    cell_text = str(cell).strip() if cell else ""
                    if cell_text:
                        for lesson in _lesson_iter(lessons):
                            atoms.append(ContentAtom(
                                atom_type="table",
                                content=cell_text,
                                source_file=rel_path,
                                location=f"sheet:{sheet_name}:r{row_idx}c{col_idx}",
                                term=term,
                                lesson=lesson,
                            ))

    return atoms


# ──────────────────────────────────────────────────────────
# Main builder
# ──────────────────────────────────────────────────────────

def _load_file_manifest():
    """Load file_manifest.json and return a path→category lookup dict.

    Files NOT in the manifest are treated as lesson_content (included in coverage).
    """
    manifest_path = BASE_DIR / "file_manifest.json"
    if not manifest_path.exists():
        return {}

    try:
        with open(manifest_path, "r", encoding="utf-8") as f:
            data = json.load(f)
    except Exception:
        return {}

    lookup = {}
    for entry in data.get("files", []):
        path = entry.get("path", "")
        category = entry.get("category", "")
        if path and category:
            # Normalize path separators for matching
            lookup[path.replace("\\", "/")] = category
    return lookup


def build_source_manifest(sources_dir, native_dir=None) -> SourceManifest:
    """Walk all source files and build a complete manifest of ContentAtoms.

    Uses file_manifest.json to exclude admin_doc, support_resource, and media
    files from coverage metrics. Files not in the manifest are assumed to be
    lesson content and are included.
    """
    sources_dir = Path(sources_dir)
    atoms = []
    source_files = []
    excluded_files = []

    # Load file manifest for inclusion/exclusion filtering
    file_manifest = _load_file_manifest()

    # Supported source extensions
    parsers = {
        ".pptx": _extract_pptx_atoms,
        ".docx": _extract_docx_atoms,
        ".pdf": _extract_pdf_atoms,
    }

    # Walk source files
    if sources_dir.exists():
        all_files = sorted(sources_dir.rglob("*"))
        for file_path in all_files:
            if not file_path.is_file() or file_path.name.startswith("~$"):
                continue
            ext = file_path.suffix.lower()
            if ext not in parsers:
                continue

            try:
                rel_path = str(file_path.relative_to(sources_dir))
            except ValueError:
                rel_path = file_path.name

            # Check file manifest — skip excluded categories
            manifest_key = rel_path.replace("\\", "/")
            category = file_manifest.get(manifest_key, "lesson_content")
            if category in _EXCLUDED_CATEGORIES:
                excluded_files.append(rel_path)
                continue

            term = extract_term_from_path(rel_path)
            lessons = extract_lesson_from_path(rel_path, term)
            if not lessons:
                lessons = [None]

            print(f"  Source: {rel_path} (term={term}, lessons={lessons})")
            source_files.append(rel_path)
            file_atoms = parsers[ext](file_path, rel_path, term, lessons)
            atoms.extend(file_atoms)

    # Walk native extracts — filter by file manifest
    if native_dir:
        native_dir = Path(native_dir)
        if native_dir.exists():
            for json_file in sorted(native_dir.rglob("*.json")):
                try:
                    rel_path = str(json_file.relative_to(native_dir))
                except ValueError:
                    rel_path = json_file.name

                # For combined extraction files, filter individual extractions inside
                # For individual files, check the file itself
                term = extract_term_from_path(rel_path)
                lessons = extract_lesson_from_path(rel_path, term)
                if not lessons:
                    lessons = [None]

                print(f"  Native: {rel_path} (term={term}, lessons={lessons})")
                source_files.append(f"native:{rel_path}")
                file_atoms = _extract_native_atoms(json_file, f"native:{rel_path}", term, lessons)

                # Filter out atoms from excluded files (admin_doc, support_resource, media)
                if file_manifest:
                    filtered_atoms = []
                    for atom in file_atoms:
                        # Check if this atom's source file is excluded
                        # Native atoms have source_file like "native:term1/Admin/File Name"
                        src = atom.source_file
                        if src.startswith("native:"):
                            src = src[7:]  # strip "native:" prefix
                        # Try matching against manifest paths
                        src_normalized = src.replace("\\", "/")
                        category = file_manifest.get(src_normalized, "")
                        if not category:
                            # Also try with common path prefixes from the manifest
                            for manifest_path, cat in file_manifest.items():
                                # Match on file name stem (handles native extracts that omit extension)
                                manifest_stem = Path(manifest_path).stem.replace("\\", "/")
                                if manifest_stem and manifest_stem in src_normalized:
                                    category = cat
                                    break
                        if category in _EXCLUDED_CATEGORIES:
                            excluded_files.append(src)
                            continue
                        filtered_atoms.append(atom)
                    atoms.extend(filtered_atoms)
                else:
                    atoms.extend(file_atoms)

    if excluded_files:
        print(f"\n  Excluded by file_manifest.json: {len(excluded_files)} files (admin/support/media)")
    print(f"  Total source atoms: {len(atoms)} from {len(source_files)} files")
    return SourceManifest(atoms=atoms, source_files=source_files, excluded_files=excluded_files)
