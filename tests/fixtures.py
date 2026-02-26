"""
Synthetic file generators for pipeline testing.
Creates controlled PPTX, PDF, and mock API response data.
"""

import io
import struct
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn


# ──────────────────────────────────────────────────────────
# PPTX generators
# ──────────────────────────────────────────────────────────

def create_test_pptx(path, slides_config):
    """Create a PPTX with controlled content.

    slides_config: list of dicts, each with:
        text: str - main text
        notes: str - speaker notes (optional)
        hyperlinks: list of {url, text} - text hyperlinks (optional)
        click_actions: list of {url, text} - click action hyperlinks (optional)
        images: list of bytes - raw image data to embed (optional)
    """
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # Blank layout

    for slide_cfg in slides_config:
        slide = prs.slides.add_slide(blank_layout)

        # Add text
        text = slide_cfg.get("text", "")
        if text:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
            tf = txBox.text_frame
            tf.text = text

        # Add hyperlinks
        for link in slide_cfg.get("hyperlinks", []):
            txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(6), Inches(1))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = link.get("text", "Click here")
            run.hyperlink.address = link["url"]

        # Add click action hyperlinks on shapes
        for action in slide_cfg.get("click_actions", []):
            txBox = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(4), Inches(1))
            tf = txBox.text_frame
            tf.text = action.get("text", "")
            txBox.click_action.hyperlink.address = action["url"]

        # Add speaker notes
        notes_text = slide_cfg.get("notes", "")
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

        # Add images
        for img_data in slide_cfg.get("images", []):
            stream = io.BytesIO(img_data)
            slide.shapes.add_picture(stream, Inches(2), Inches(2), Inches(3), Inches(2))

    prs.save(str(path))


def create_minimal_png(width=2, height=2, color=(255, 0, 0)):
    """Create a minimal valid PNG file in memory."""
    import zlib

    def _chunk(chunk_type, data):
        c = chunk_type + data
        crc = struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)
        return struct.pack(">I", len(data)) + c + crc

    # IHDR
    ihdr_data = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    # IDAT - raw image data
    raw = b""
    for _ in range(height):
        raw += b"\x00"  # filter byte
        for _ in range(width):
            raw += bytes(color)
    compressed = zlib.compress(raw)

    png = b"\x89PNG\r\n\x1a\n"
    png += _chunk(b"IHDR", ihdr_data)
    png += _chunk(b"IDAT", compressed)
    png += _chunk(b"IEND", b"")
    return png


def create_corrupted_pptx(path):
    """Create a file that looks like a PPTX but is actually corrupted."""
    with open(path, "wb") as f:
        f.write(b"PK\x03\x04" + b"\x00" * 100)  # Broken ZIP header


def create_zero_byte_file(path):
    """Create a zero-byte file."""
    with open(path, "wb") as f:
        pass


# ──────────────────────────────────────────────────────────
# PDF generators
# ──────────────────────────────────────────────────────────

def create_test_pdf(path, pages_config):
    """Create a PDF with controlled content using PyPDF2.

    pages_config: list of dicts, each with:
        text: str - page text
        links: list of {url} - hyperlinks (optional)
    """
    from PyPDF2 import PdfWriter
    from PyPDF2.generic import (
        ArrayObject, DictionaryObject, NameObject,
        TextStringObject, NumberObject, FloatObject,
    )

    writer = PdfWriter()

    for page_cfg in pages_config:
        # Create a blank page
        writer.add_blank_page(width=612, height=792)
        page = writer.pages[-1]

        # Add link annotations
        for link in page_cfg.get("links", []):
            annot = DictionaryObject()
            annot.update({
                NameObject("/Type"): NameObject("/Annot"),
                NameObject("/Subtype"): NameObject("/Link"),
                NameObject("/Rect"): ArrayObject([
                    FloatObject(72), FloatObject(700),
                    FloatObject(200), FloatObject(720),
                ]),
                NameObject("/A"): DictionaryObject({
                    NameObject("/Type"): NameObject("/Action"),
                    NameObject("/S"): NameObject("/URI"),
                    NameObject("/URI"): TextStringObject(link["url"]),
                }),
            })

            if "/Annots" not in page:
                page[NameObject("/Annots")] = ArrayObject()
            page[NameObject("/Annots")].append(annot)

    with open(path, "wb") as f:
        writer.write(f)


def create_empty_pdf(path):
    """Create a valid PDF with zero content pages."""
    from PyPDF2 import PdfWriter
    writer = PdfWriter()
    with open(path, "wb") as f:
        writer.write(f)


def create_encrypted_pdf(path):
    """Create a password-protected PDF."""
    from PyPDF2 import PdfWriter
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    writer.encrypt("secretpassword")
    with open(path, "wb") as f:
        writer.write(f)


def create_pdf_with_many_links(path, count=100):
    """Create a PDF with many link annotations."""
    pages = []
    links_per_page = 25
    for page_num in range(0, count, links_per_page):
        batch = min(links_per_page, count - page_num)
        links = [{"url": f"https://example.com/link{page_num + i}"} for i in range(batch)]
        pages.append({"text": f"Page with links", "links": links})
    create_test_pdf(path, pages)


# ──────────────────────────────────────────────────────────
# Mock Google API responses
# ──────────────────────────────────────────────────────────

def build_slides_api_response(slides_config):
    """Build a dict matching Google Slides API presentations().get() response.

    slides_config: list of dicts, each with:
        texts: list of str - text content
        links: list of {url, text} - hyperlinks in text runs
        videos: list of {url, source, video_id} - embedded videos
        table: list of list of str - table rows (optional)
        notes: str - speaker notes (optional)
        notes_links: list of {url, text} - links in speaker notes (optional)
    """
    slides = []
    for i, cfg in enumerate(slides_config, 1):
        page_elements = []

        # Text shapes with links
        for j, text in enumerate(cfg.get("texts", [])):
            text_elements = [{
                "textRun": {
                    "content": text,
                    "style": {},
                },
            }]
            page_elements.append({
                "objectId": f"slide{i}_shape{j}",
                "shape": {
                    "text": {"textElements": text_elements},
                },
            })

        # Text shapes with hyperlinks
        for link in cfg.get("links", []):
            text_elements = [{
                "textRun": {
                    "content": link.get("text", ""),
                    "style": {
                        "link": {"url": link["url"]},
                    },
                },
            }]
            page_elements.append({
                "objectId": f"slide{i}_link_{link['url'][:10]}",
                "shape": {
                    "text": {"textElements": text_elements},
                },
            })

        # Embedded videos
        for vid in cfg.get("videos", []):
            page_elements.append({
                "objectId": f"slide{i}_video",
                "video": {
                    "url": vid.get("url", ""),
                    "source": vid.get("source", "YOUTUBE"),
                    "id": vid.get("video_id", ""),
                },
            })

        # Tables with optional links in cells
        table_data = cfg.get("table")
        if table_data:
            table_rows = []
            for row in table_data:
                table_cells = []
                for cell_text in row:
                    cell_te = [{"textRun": {"content": cell_text, "style": {}}}]
                    table_cells.append({"text": {"textElements": cell_te}})
                table_rows.append({"tableCells": table_cells})
            page_elements.append({
                "objectId": f"slide{i}_table",
                "table": {"tableRows": table_rows},
            })

        # Tables with links
        table_links = cfg.get("table_links")
        if table_links:
            table_rows = []
            for row in table_links:
                table_cells = []
                for cell in row:
                    if isinstance(cell, dict):
                        cell_te = [{
                            "textRun": {
                                "content": cell.get("text", ""),
                                "style": {"link": {"url": cell["url"]}},
                            },
                        }]
                    else:
                        cell_te = [{"textRun": {"content": cell, "style": {}}}]
                    table_cells.append({"text": {"textElements": cell_te}})
                table_rows.append({"tableCells": table_cells})
            page_elements.append({
                "objectId": f"slide{i}_linked_table",
                "table": {"tableRows": table_rows},
            })

        # Speaker notes
        notes_elements = []
        notes_text = cfg.get("notes", "")
        if notes_text:
            notes_elements.append({
                "textRun": {"content": notes_text, "style": {}},
            })
        for nl in cfg.get("notes_links", []):
            notes_elements.append({
                "textRun": {
                    "content": nl.get("text", ""),
                    "style": {"link": {"url": nl["url"]}},
                },
            })

        notes_page = {}
        if notes_elements:
            notes_page = {
                "pageElements": [{
                    "shape": {
                        "shapeType": "TEXT_BOX",
                        "text": {"textElements": notes_elements},
                    },
                }],
            }

        slide_obj = {
            "objectId": f"slide_{i}",
            "pageElements": page_elements,
            "slideProperties": {},
        }
        if notes_page:
            slide_obj["slideProperties"]["notesPage"] = notes_page

        slides.append(slide_obj)

    return {
        "presentationId": "test_presentation_id",
        "title": "Test Presentation",
        "slides": slides,
    }


def build_docs_api_response(blocks_config):
    """Build a dict matching Google Docs API documents().get() response.

    blocks_config: list of dicts, each with:
        text: str - paragraph text
        style: str - paragraph style (HEADING_1, HEADING_2, HEADING_3, NORMAL_TEXT)
        links: list of {url, text} - hyperlinks in text (optional)
    """
    content = []
    for block in blocks_config:
        elements = []

        # Plain text
        text = block.get("text", "")
        if text:
            elements.append({
                "textRun": {
                    "content": text,
                    "textStyle": {},
                },
            })

        # Links
        for link in block.get("links", []):
            elements.append({
                "textRun": {
                    "content": link.get("text", ""),
                    "textStyle": {
                        "link": {"url": link["url"]},
                    },
                },
            })

        paragraph = {
            "paragraph": {
                "paragraphStyle": {
                    "namedStyleType": block.get("style", "NORMAL_TEXT"),
                },
                "elements": elements,
            },
        }
        content.append(paragraph)

    # Table blocks
    for block in blocks_config:
        table = block.get("table")
        if table:
            table_rows = []
            for row in table:
                table_cells = []
                for cell_text in row:
                    if isinstance(cell_text, dict):
                        cell_elements = [{
                            "textRun": {
                                "content": cell_text.get("text", ""),
                                "textStyle": {
                                    "link": {"url": cell_text["url"]},
                                },
                            },
                        }]
                    else:
                        cell_elements = [{
                            "textRun": {
                                "content": cell_text,
                                "textStyle": {},
                            },
                        }]
                    table_cells.append({
                        "content": [{
                            "paragraph": {"elements": cell_elements},
                        }],
                    })
                table_rows.append({"tableCells": table_cells})

            content.append({"table": {"tableRows": table_rows}})

    return {
        "documentId": "test_doc_id",
        "title": "Test Document",
        "body": {"content": content},
    }


def create_video_file(path, size_kb=10):
    """Create a fake video file (just random bytes with .mp4 extension)."""
    with open(path, "wb") as f:
        f.write(b"\x00" * (size_kb * 1024))
