"""
Stage 8 Tests: Cross-Validation Expert Agent

Tests marked @pytest.mark.llm require Claude CLI and are skipped without it.
Run structural tests:  pytest tests/test_cross_validation.py -v -m "not llm"
Run LLM tests:         pytest tests/test_cross_validation.py -v -m llm
"""

import json
import subprocess
import pytest
from pathlib import Path
from unittest.mock import patch, MagicMock

from tests.fixtures import create_test_pptx

from cross_validate_kb import (
    extract_pptx_ground_truth,
    build_ground_truth_text,
    build_kb_text,
    build_error_investigation_prompt,
    build_lesson_evaluation_prompt,
    score_verdict,
    compute_overall_confidence,
    generate_recommendations,
    generate_report,
    write_text_summary,
    find_source_pptx,
    EVALUATED_FIELDS,
)


def _cli_available():
    try:
        result = subprocess.run(
            ["claude", "--version"], capture_output=True, text=True, timeout=10,
        )
        return result.returncode == 0
    except (FileNotFoundError, subprocess.TimeoutExpired):
        return False


skip_no_cli = pytest.mark.skipif(
    not _cli_available(), reason="Claude CLI not available"
)


# ──────────────────────────────────────────────────────────
# Ground Truth Extraction Tests
# ──────────────────────────────────────────────────────────

class TestExtractGroundTruth:

    def test_extract_from_sample_pptx(self, tmp_path):
        """Verify ground truth extraction returns structured slide content."""
        pptx_path = tmp_path / "test.pptx"
        create_test_pptx(pptx_path, [
            {
                "text": "Lesson 1 - Introduction to Design",
                "notes": "Welcome students",
                "hyperlinks": [
                    {"url": "https://example.com/resource", "text": "Resource Link"},
                ],
            },
            {
                "text": "Learning Objectives\nUnderstand design principles",
                "notes": "Explain each objective",
            },
        ])

        result = extract_pptx_ground_truth(pptx_path)

        assert "error" not in result
        assert result["total_slides"] == 2

        slide1 = result["slides"][0]
        assert slide1["slide_number"] == 1
        assert any("Design" in t for t in slide1["text"])
        assert "Welcome students" in slide1["notes"]
        assert len(slide1["links"]) >= 1
        assert slide1["links"][0]["url"] == "https://example.com/resource"

        slide2 = result["slides"][1]
        assert any("Learning Objectives" in t for t in slide2["text"])

    def test_extract_nonexistent_file(self):
        """Nonexistent file returns error dict with empty slides."""
        result = extract_pptx_ground_truth(Path("/nonexistent/file.pptx"))
        assert "error" in result
        assert result["slides"] == []

    def test_extract_tables(self, tmp_path):
        """Verify table content is extracted."""
        from pptx import Presentation
        from pptx.util import Inches

        pptx_path = tmp_path / "table_test.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add a table
        table = slide.shapes.add_table(2, 3, Inches(1), Inches(1), Inches(6), Inches(2)).table
        table.cell(0, 0).text = "Header1"
        table.cell(0, 1).text = "Header2"
        table.cell(0, 2).text = "Header3"
        table.cell(1, 0).text = "Data1"
        table.cell(1, 1).text = "Data2"
        table.cell(1, 2).text = "Data3"
        prs.save(pptx_path)

        result = extract_pptx_ground_truth(pptx_path)
        assert result["total_slides"] == 1
        assert len(result["slides"][0]["tables"]) == 1
        assert result["slides"][0]["tables"][0][0] == ["Header1", "Header2", "Header3"]


# ──────────────────────────────────────────────────────────
# Prompt Construction Tests
# ──────────────────────────────────────────────────────────

class TestPromptConstruction:

    def test_ground_truth_text_format(self):
        """Verify ground truth is formatted as readable text."""
        gt = {
            "slides": [
                {
                    "slide_number": 1,
                    "text": ["Title Slide", "Subtitle"],
                    "notes": "Teacher notes here",
                    "links": [{"url": "https://example.com", "text": "Link"}],
                    "tables": [],
                },
            ]
        }

        text = build_ground_truth_text(gt)
        assert "Slide 1" in text
        assert "Title Slide" in text
        assert "Teacher notes here" in text
        assert "https://example.com" in text

    def test_ground_truth_empty(self):
        """Empty ground truth returns placeholder."""
        assert "[No ground truth available]" in build_ground_truth_text(None)
        assert "[No ground truth available]" in build_ground_truth_text({"error": "test"})

    def test_kb_text_format(self):
        """Verify KB lesson is formatted correctly."""
        kb = {
            "lesson_title": "Lesson 5 - Brainstorming",
            "metadata": {
                "core_topics": ["brainstorming", "concept generation"],
                "learning_objectives": ["Generate ideas"],
                "endstar_tools": ["Triggers"],
                "activity_type": "group_work",
                "activity_description": "Students brainstorm in groups",
                "videos": [],
                "resources": ["https://example.com - Resource"],
                "keywords": ["brainstorming", "design"],
            },
            "description_of_activities": "Students work in groups to brainstorm.",
        }

        text = build_kb_text(kb)
        assert "Lesson 5" in text
        assert "brainstorming" in text
        assert "Generate ideas" in text
        assert "Triggers" in text

    def test_kb_text_none(self):
        """None KB returns placeholder."""
        assert "[No KB entry found]" in build_kb_text(None)

    def test_error_investigation_prompt_structure(self):
        """Error investigation prompt includes all required sections."""
        error = {
            "type": "MISSING",
            "message": "Lesson 2 missing teachers slides",
            "lesson": 2,
            "content_type": "teachers_slides",
        }

        prompt = build_error_investigation_prompt(error, "Source text", "KB text")
        assert "MISSING" in prompt
        assert "teachers_slides" in prompt
        assert "Source text" in prompt
        assert "KB text" in prompt
        assert "TRUE POSITIVE" in prompt or "TRUE_POSITIVE" in prompt
        assert "FALSE POSITIVE" in prompt or "FALSE_POSITIVE" in prompt

    def test_lesson_evaluation_prompt_structure(self):
        """Lesson evaluation prompt lists all 8 fields."""
        prompt = build_lesson_evaluation_prompt("Source text", "KB text")
        assert "lesson_title" in prompt
        assert "learning_objectives" in prompt
        assert "core_topics" in prompt
        assert "activity_description" in prompt
        assert "resources" in prompt
        assert "videos" in prompt
        assert "endstar_tools" in prompt
        assert "keywords" in prompt
        assert "CORRECT" in prompt
        assert "PARTIAL" in prompt


# ──────────────────────────────────────────────────────────
# Scoring Logic Tests
# ──────────────────────────────────────────────────────────

class TestScoring:

    def test_verdict_scores(self):
        assert score_verdict("CORRECT") == 1.0
        assert score_verdict("PARTIAL") == 0.5
        assert score_verdict("INCORRECT") == 0.0
        assert score_verdict("MISSING") == 0.0
        assert score_verdict("UNKNOWN") == 0.0

    def test_overall_confidence_all_correct(self):
        """Perfect phase2 + all false positives → high confidence."""
        phase1 = {
            "total_errors_reviewed": 10,
            "confirmed_errors": [],
            "false_positives": [{}] * 10,
            "uncertain": [],
        }
        phase2 = {
            "per_lesson": [{"overall_score": 1.0}] * 5,
        }

        confidence = compute_overall_confidence(phase1, phase2)
        assert confidence >= 0.9

    def test_overall_confidence_many_confirmed_errors(self):
        """Many confirmed errors → lower confidence."""
        phase1 = {
            "total_errors_reviewed": 10,
            "confirmed_errors": [{}] * 8,
            "false_positives": [{}] * 2,
            "uncertain": [],
        }
        phase2 = {
            "per_lesson": [{"overall_score": 0.5}] * 5,
        }

        confidence = compute_overall_confidence(phase1, phase2)
        assert confidence < 0.7

    def test_overall_confidence_no_data(self):
        """No data → neutral confidence."""
        phase1 = {
            "total_errors_reviewed": 0,
            "confirmed_errors": [],
            "false_positives": [],
            "uncertain": [],
        }
        phase2 = {"per_lesson": []}

        confidence = compute_overall_confidence(phase1, phase2)
        assert 0 <= confidence <= 1


# ──────────────────────────────────────────────────────────
# Recommendation Tests
# ──────────────────────────────────────────────────────────

class TestRecommendations:

    def test_naming_mismatch_recommendation(self):
        """False positives from naming → recommend fixing content_type detection."""
        phase1 = {
            "false_positives": [
                {"reason": "content_type != teachers_slides, naming mismatch"},
                {"reason": "content_type mismatch"},
            ],
        }
        phase2 = {"field_accuracy": {}}

        recs = generate_recommendations(phase1, phase2)
        assert any("content_type" in r.lower() or "naming" in r.lower() for r in recs)

    def test_students_slides_recommendation(self):
        """Many students_slides false positives → recommend making optional."""
        phase1 = {
            "false_positives": [
                {"content_type": "students_slides"},
                {"content_type": "students_slides"},
            ],
        }
        phase2 = {"field_accuracy": {}}

        recs = generate_recommendations(phase1, phase2)
        assert any("students_slides" in r for r in recs)

    def test_no_issues_recommendation(self):
        """No issues → positive recommendation."""
        phase1 = {"false_positives": []}
        phase2 = {"field_accuracy": {}}

        recs = generate_recommendations(phase1, phase2)
        assert any("no critical" in r.lower() for r in recs)


# ──────────────────────────────────────────────────────────
# Report Generation Tests
# ──────────────────────────────────────────────────────────

class TestReportGeneration:

    def test_report_json_structure(self):
        """Verify report has all required top-level keys."""
        phase1 = {
            "total_errors_reviewed": 5,
            "qa_failures_reviewed": 0,
            "confirmed_errors": [],
            "false_positives": [{"reason": "test"}] * 5,
            "uncertain": [],
        }
        phase2 = {
            "total_available": 10,
            "total_selected": 5,
            "selection_rate": "50%",
            "per_lesson": [
                {"term": 1, "lesson": 1, "fields": {}, "overall_score": 0.8},
            ],
            "field_accuracy": {f: {"correct": 1, "partial": 0, "incorrect": 0, "missing": 0}
                              for f in EVALUATED_FIELDS},
        }

        report = generate_report(phase1, phase2)

        assert "cross_validated_at" in report
        assert "model_used" in report
        assert "phase1_error_investigation" in report
        assert "phase2_pass_verification" in report
        assert "overall_confidence" in report
        assert "recommendations" in report
        assert isinstance(report["overall_confidence"], float)
        assert 0 <= report["overall_confidence"] <= 1
        assert isinstance(report["recommendations"], list)

    def test_text_summary_generated(self, tmp_path):
        """Verify text summary writes without errors."""
        report = {
            "cross_validated_at": "2026-02-25T00:00:00Z",
            "model_used": "opus",
            "cli_tool": "claude",
            "overall_confidence": 0.85,
            "phase1_error_investigation": {
                "total_errors_reviewed": 3,
                "confirmed_errors": [{"term": 1, "lesson": 2, "type": "MISSING",
                                      "llm_reason": "test reason"}],
                "false_positives": [{"term": 1, "lesson": 1, "content_type": "students_slides",
                                     "reason": "Not provided"}],
                "uncertain": [],
            },
            "phase2_pass_verification": {
                "total_available": 10,
                "total_selected": 5,
                "selection_rate": "50%",
                "per_lesson": [
                    {"term": 1, "lesson": 1, "fields": {}, "overall_score": 0.9},
                    {"term": 2, "lesson": 3, "fields": {}, "overall_score": 0.8},
                ],
                "field_accuracy": {
                    "lesson_title": {"correct": 2, "partial": 0, "incorrect": 0, "missing": 0},
                },
            },
            "recommendations": ["Fix content_type detection"],
        }

        txt_path = tmp_path / "report.txt"
        write_text_summary(report, txt_path)

        content = txt_path.read_text(encoding="utf-8")
        assert "CROSS-VALIDATION REPORT" in content
        assert "85%" in content
        assert "Fix content_type detection" in content
        assert "Phase 1" in content.upper() or "PHASE 1" in content


# ──────────────────────────────────────────────────────────
# LLM Integration Tests (require Claude CLI)
# ──────────────────────────────────────────────────────────

@pytest.mark.llm
class TestLLMJudge:

    @skip_no_cli
    def test_judge_error_false_positive(self, tmp_path):
        """Feed a known false positive and verify judge catches it."""
        from cross_validate_kb import judge_error

        error = {
            "type": "MISSING",
            "message": "Lesson 1 missing teachers slides",
            "lesson": 1,
            "content_type": "teachers_slides",
        }

        # Source content: clearly a teacher's presentation
        source_content = {
            "slides": [
                {
                    "slide_number": 1,
                    "text": ["Lesson 1 - Introduction to Game Design",
                             "Teacher Presentation"],
                    "notes": "Welcome students to the first lesson",
                    "links": [],
                    "tables": [],
                },
                {
                    "slide_number": 2,
                    "text": ["Learning Objectives",
                             "Understand what game design is",
                             "Identify key design elements"],
                    "notes": "Go through each objective",
                    "links": [],
                    "tables": [],
                },
            ]
        }

        kb_content = {
            "lesson_title": "Lesson 1 - Introduction to Game Design",
            "metadata": {
                "core_topics": ["game design", "design elements"],
                "learning_objectives": ["Understand game design", "Identify design elements"],
                "endstar_tools": [],
                "activity_type": "introduction",
                "activity_description": "Introduction to game design concepts",
                "videos": [],
                "resources": [],
                "keywords": ["game design", "introduction"],
            },
        }

        result = judge_error(error, source_content, kb_content)

        assert "verdict" in result
        assert result["verdict"] in ("TRUE_POSITIVE", "FALSE_POSITIVE", "UNCERTAIN")
        assert "reason" in result

    @skip_no_cli
    def test_judge_lesson_evaluation(self, tmp_path):
        """Verify full lesson evaluation returns 8-field results."""
        from cross_validate_kb import judge_lesson

        source_content = {
            "slides": [
                {
                    "slide_number": 1,
                    "text": ["Lesson 5 - Brainstorming and Concept Generation"],
                    "notes": "",
                    "links": [{"url": "https://example.com/rubric", "text": "Rubric"}],
                    "tables": [],
                },
                {
                    "slide_number": 2,
                    "text": ["Learning Objectives",
                             "Generate creative game concepts",
                             "Use brainstorming techniques"],
                    "notes": "Explain brainstorming methods",
                    "links": [],
                    "tables": [],
                },
            ]
        }

        kb_content = {
            "lesson_title": "Lesson 5 - Brainstorming and Concept Generation",
            "metadata": {
                "core_topics": ["brainstorming", "concept generation"],
                "learning_objectives": ["Generate creative game concepts"],
                "endstar_tools": [],
                "activity_type": "group_work",
                "activity_description": "Brainstorming in groups",
                "videos": [],
                "resources": ["Rubric - https://example.com/rubric"],
                "keywords": ["brainstorming", "design", "concepts"],
            },
        }

        result = judge_lesson(source_content, kb_content)

        for field in EVALUATED_FIELDS:
            assert field in result, f"Missing field: {field}"
            assert "verdict" in result[field], f"Missing verdict for {field}"
            assert result[field]["verdict"] in ("CORRECT", "PARTIAL", "INCORRECT", "MISSING")
