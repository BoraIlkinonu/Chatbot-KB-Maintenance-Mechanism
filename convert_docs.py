"""
Stage 2: Document Conversion
Converts PPTX, DOCX, XLSX, PDF to structured Markdown/CSV.
Preserves tables as structured data. Extracts speaker notes from PPTX.
"""

import sys
import os
import json
import csv
import traceback
from pathlib import Path
from datetime import datetime, timezone

sys.stdout.reconfigure(encoding="utf-8")

from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

from config import SOURCES_DIR, CONVERTED_DIR, LOGS_DIR, MEDIA_DIR


def _iter_all_shapes(shapes):
    """Recursively yield all shapes, descending into GroupShapes."""
    try:
        from pptx.shapes.group import GroupShape
    except ImportError:
        GroupShape = type(None)
    for shape in shapes:
        yield shape
        if isinstance(shape, GroupShape):
            yield from _iter_all_shapes(shape.shapes)


def convert_docx(docx_path, output_dir):
    """Convert DOCX to Markdown with table preservation."""
    try:
        doc = Document(docx_path)
        output_file = output_dir / (docx_path.stem + ".md")

        content = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            style_name = para.style.name.lower() if para.style else ""
            if "heading 1" in style_name or "title" in style_name:
                content.append(f"# {text}\n")
            elif "heading 2" in style_name:
                content.append(f"## {text}\n")
            elif "heading 3" in style_name:
                content.append(f"### {text}\n")
            elif "list" in style_name:
                content.append(f"- {text}\n")
            else:
                content.append(f"{text}\n")

        # Extract tables as structured markdown
        for table in doc.tables:
            if not table.rows:
                continue
            headers = [cell.text.strip().replace("\n", " ") for cell in table.rows[0].cells]
            content.append("\n| " + " | ".join(headers) + " |")
            content.append("| " + " | ".join(["---"] * len(headers)) + " |")
            for row in table.rows[1:]:
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                content.append("| " + " | ".join(cells) + " |")
            content.append("")

        with open(output_file, "w", encoding="utf-8") as f:
            f.write("\n".join(content))

        return str(output_file), True, None
    except Exception as e:
        return str(docx_path), False, str(e)


def convert_pptx(pptx_path, output_dir):
    """Convert PPTX to Markdown. Includes speaker notes and table data."""
    try:
        prs = Presentation(pptx_path)
        output_file = output_dir / (pptx_path.stem + ".md")

        content = [f"# {pptx_path.stem}\n"]

        for slide_num, slide in enumerate(prs.slides, 1):
            content.append(f"\n## Slide {slide_num}\n")

            for shape in _iter_all_shapes(slide.shapes):
                # Text content
                if hasattr(shape, "text") and shape.text.strip():
                    content.append(f"{shape.text.strip()}\n")

                # Tables
                if shape.has_table:
                    table = shape.table
                    if table.rows:
                        headers = [
                            cell.text.strip().replace("\n", " ")
                            for cell in table.rows[0].cells
                        ]
                        content.append("\n| " + " | ".join(headers) + " |")
                        content.append("| " + " | ".join(["---"] * len(headers)) + " |")
                        for row in list(table.rows)[1:]:
                            cells = [
                                cell.text.strip().replace("\n", " ")
                                for cell in row.cells
                            ]
                            content.append("| " + " | ".join(cells) + " |")
                        content.append("")

            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    content.append(f"\n**Speaker Notes:**\n{notes}\n")

            content.append("\n---\n")

        with open(output_file, "w", encoding="utf-8") as f:
            f.write("\n".join(content))

        return str(output_file), True, None
    except Exception as e:
        return str(pptx_path), False, str(e)


def convert_xlsx(xlsx_path, output_dir):
    """Convert XLSX sheets to CSV files."""
    try:
        wb = load_workbook(xlsx_path, data_only=True)
        output_files = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            clean = "".join(c if c.isalnum() or c in " -_" else "_" for c in sheet_name)
            out_file = output_dir / f"{xlsx_path.stem}_{clean}.csv"

            with open(out_file, "w", newline="", encoding="utf-8") as f:
                writer = csv.writer(f)
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        writer.writerow(
                            [str(cell) if cell is not None else "" for cell in row]
                        )
            output_files.append(str(out_file))

        return output_files, True, None
    except Exception as e:
        return str(xlsx_path), False, str(e)


def extract_pdf_images(reader, pdf_path):
    """Extract images from PDF via XObject parsing.
    Saves to media/pdf_images/<stem>/. Returns list of image metadata."""
    images = []
    pdf_img_dir = MEDIA_DIR / "pdf_images" / pdf_path.stem.replace(" ", "_")
    pdf_img_dir.mkdir(parents=True, exist_ok=True)

    img_index = 0
    for page_num, page in enumerate(reader.pages, 1):
        resources = page.get("/Resources")
        if not resources:
            continue
        xobjects = resources.get("/XObject")
        if not xobjects:
            continue

        xobj_dict = xobjects.get_object() if hasattr(xobjects, "get_object") else xobjects
        if not isinstance(xobj_dict, dict):
            continue

        for obj_name in xobj_dict:
            try:
                xobj = xobj_dict[obj_name]
                obj = xobj.get_object() if hasattr(xobj, "get_object") else xobj
                if not hasattr(obj, "get"):
                    continue
                subtype = obj.get("/Subtype", "")
                if str(subtype) != "/Image":
                    continue

                img_index += 1
                # Determine extension from filter
                img_filter = obj.get("/Filter", "")
                filter_str = str(img_filter)
                if "DCTDecode" in filter_str:
                    ext = ".jpg"
                elif "FlateDecode" in filter_str:
                    ext = ".png"
                elif "JPXDecode" in filter_str:
                    ext = ".jp2"
                else:
                    ext = ".bin"

                img_data = obj.get_data() if hasattr(obj, "get_data") else obj.getData()
                output_name = f"image_{img_index:03d}{ext}"
                output_path = pdf_img_dir / output_name

                with open(output_path, "wb") as f:
                    f.write(img_data)

                images.append({
                    "image_path": str(output_path),
                    "page_number": page_num,
                    "index": img_index,
                    "size_bytes": len(img_data),
                    "extension": ext,
                    "source": "pdf",
                })
            except Exception:
                continue

    return images


def extract_pdf_links(reader):
    """Extract hyperlinks from PDF annotations.
    Returns list of {url, page_number}."""
    links = []
    for page_num, page in enumerate(reader.pages, 1):
        annots = page.get("/Annots")
        if not annots:
            continue

        annot_list = annots.get_object() if hasattr(annots, "get_object") else annots
        if not isinstance(annot_list, (list, tuple)):
            continue

        for annot_ref in annot_list:
            try:
                annot = annot_ref.get_object() if hasattr(annot_ref, "get_object") else annot_ref
                if not isinstance(annot, dict):
                    continue
                subtype = str(annot.get("/Subtype", ""))
                if subtype != "/Link":
                    continue

                action = annot.get("/A")
                if action:
                    action_obj = action.get_object() if hasattr(action, "get_object") else action
                    if isinstance(action_obj, dict):
                        uri = action_obj.get("/URI", "")
                        if uri:
                            uri_str = str(uri)
                            links.append({
                                "url": uri_str,
                                "page_number": page_num,
                            })
            except Exception:
                continue

    return links


def convert_pdf(pdf_path, output_dir):
    """Convert PDF to Markdown (best effort). Also extracts images and links.
    Uses pypdf (preferred) or PyPDF2 as fallback."""
    try:
        # Prefer pypdf (maintained successor), fall back to PyPDF2
        try:
            from pypdf import PdfReader
        except ImportError:
            from PyPDF2 import PdfReader

        reader = PdfReader(pdf_path)
        output_file = output_dir / (pdf_path.stem + ".md")

        content = [f"# {pdf_path.stem}\n"]
        for page_num, page in enumerate(reader.pages, 1):
            try:
                text = page.extract_text()
                if text:
                    content.append(f"\n## Page {page_num}\n")
                    content.append(text + "\n")
                    content.append("---\n")
            except Exception:
                content.append(f"\n## Page {page_num}\n")
                content.append("[Page could not be extracted]\n")
                content.append("---\n")

        with open(output_file, "w", encoding="utf-8") as f:
            f.write("\n".join(content))

        # Extract images and links (best effort)
        try:
            pdf_images = extract_pdf_images(reader, pdf_path)
        except Exception:
            pdf_images = []
        try:
            pdf_links = extract_pdf_links(reader)
        except Exception:
            pdf_links = []

        return str(output_file), True, None, {"images": pdf_images, "links": pdf_links}
    except ImportError:
        return str(pdf_path), False, "pypdf/PyPDF2 not installed", {}
    except Exception as e:
        return str(pdf_path), False, str(e), {}


def run_conversion(source_dir=None):
    """Run document conversion on all supported files in sources."""
    print("=" * 60)
    print("  Stage 2: Document Conversion")
    print("=" * 60)
    print()

    src = Path(source_dir) if source_dir else SOURCES_DIR
    CONVERTED_DIR.mkdir(parents=True, exist_ok=True)

    results = {
        "converted_at": datetime.now(timezone.utc).isoformat(),
        "files": [],
        "summary": {"success": 0, "failed": 0, "by_type": {}},
        "pdf_metadata": {
            "total_images": 0,
            "total_links": 0,
            "files": [],
        },
    }

    converters = {
        ".docx": ("DOCX", convert_docx),
        ".pptx": ("PPTX", convert_pptx),
        ".xlsx": ("XLSX", convert_xlsx),
        ".pdf":  ("PDF",  convert_pdf),
    }

    all_files = list(src.rglob("*"))
    all_files = [f for f in all_files if f.is_file() and not f.name.startswith("~$")]

    for file_path in all_files:
        ext = file_path.suffix.lower()
        if ext not in converters:
            continue

        label, converter = converters[ext]

        # Mirror directory structure in output
        try:
            rel = file_path.relative_to(src)
        except ValueError:
            rel = Path(file_path.name)
        out_dir = CONVERTED_DIR / rel.parent
        out_dir.mkdir(parents=True, exist_ok=True)

        print(f"[{label}] {file_path.name}")

        # PDF converter returns 4 values; others return 3
        if label == "PDF":
            output, success, error, pdf_extra = converter(file_path, out_dir)
            if pdf_extra:
                pdf_images = pdf_extra.get("images", [])
                pdf_links = pdf_extra.get("links", [])
                if pdf_images or pdf_links:
                    print(f"  PDF extras: {len(pdf_images)} images, {len(pdf_links)} links")
                    results["pdf_metadata"]["total_images"] += len(pdf_images)
                    results["pdf_metadata"]["total_links"] += len(pdf_links)
                    results["pdf_metadata"]["files"].append({
                        "source": str(file_path),
                        "relative_path": str(rel),
                        "images": pdf_images,
                        "links": pdf_links,
                    })
        else:
            output, success, error = converter(file_path, out_dir)

        results["files"].append({
            "source": str(file_path),
            "type": label,
            "output": output,
            "success": success,
            "error": error,
        })

        if success:
            results["summary"]["success"] += 1
        else:
            results["summary"]["failed"] += 1
            print(f"  ERROR: {error}")

        results["summary"]["by_type"][label] = (
            results["summary"]["by_type"].get(label, 0) + 1
        )

    # Save conversion log
    LOGS_DIR.mkdir(parents=True, exist_ok=True)
    log_path = LOGS_DIR / f"conversion_{datetime.now(timezone.utc).strftime('%Y%m%d_%H%M%S')}.json"
    with open(log_path, "w", encoding="utf-8") as f:
        json.dump(results, f, indent=2, ensure_ascii=False)

    # Save PDF metadata separately for consolidation
    if results["pdf_metadata"]["files"]:
        MEDIA_DIR.mkdir(parents=True, exist_ok=True)
        pdf_meta_path = MEDIA_DIR / "pdf_extraction_metadata.json"
        with open(pdf_meta_path, "w", encoding="utf-8") as f:
            json.dump(results["pdf_metadata"], f, indent=2, ensure_ascii=False)
        print(f"PDF metadata saved: {pdf_meta_path}")

    print(f"\nConversion: {results['summary']['success']} success, {results['summary']['failed']} failed")
    if results["pdf_metadata"]["total_images"] > 0 or results["pdf_metadata"]["total_links"] > 0:
        print(f"PDF images: {results['pdf_metadata']['total_images']}, PDF links: {results['pdf_metadata']['total_links']}")
    print(f"Log saved: {log_path}")
    print("=" * 60)

    return results


if __name__ == "__main__":
    run_conversion()
