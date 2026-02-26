"""
Media Extraction Pipeline for Term 2 Teacher Resources
Extracts images from PPTX files and keyframes from videos
WITH SLIDE NUMBER TRACKING
"""

import os
import zipfile
import shutil
import subprocess
import json
import re
import xml.etree.ElementTree as ET
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io

# Configuration
SOURCE_DIR = Path(r"D:\Term 3 QA\Teacher Resources - Term 2\Teacher Resources")
OUTPUT_DIR = Path(r"D:\Term 3 QA\Teacher Resources - Term 2\Extracted Media")
KEYFRAME_INTERVAL = 10  # Extract keyframe every N seconds

# XML namespaces used in PPTX files
NAMESPACES = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'rel': 'http://schemas.openxmlformats.org/package/2006/relationships'
}

def setup_output_dirs():
    """Create output directory structure"""
    dirs = [
        OUTPUT_DIR / "pptx_images",
        OUTPUT_DIR / "video_keyframes",
        OUTPUT_DIR / "video_transcripts",
        OUTPUT_DIR / "metadata"
    ]
    for d in dirs:
        d.mkdir(parents=True, exist_ok=True)
    return dirs

def build_slide_image_mapping(zip_ref):
    """
    Parse PPTX XML to build a mapping of which images appear on which slides.
    Returns: dict mapping media filename -> list of slide numbers
    """
    media_to_slides = {}

    # Find all slide files
    slide_files = sorted([f for f in zip_ref.namelist() if re.match(r'ppt/slides/slide\d+\.xml$', f)])

    for slide_file in slide_files:
        # Extract slide number from filename (e.g., "ppt/slides/slide1.xml" -> 1)
        slide_num = int(re.search(r'slide(\d+)\.xml', slide_file).group(1))

        # Get the relationship file for this slide
        rels_file = slide_file.replace('slides/', 'slides/_rels/') + '.rels'

        if rels_file not in zip_ref.namelist():
            continue

        # Parse the relationship file to get rId -> media file mapping
        rels_content = zip_ref.read(rels_file).decode('utf-8')
        rels_root = ET.fromstring(rels_content)

        rid_to_media = {}
        for rel in rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
            rel_id = rel.get('Id')
            target = rel.get('Target')
            rel_type = rel.get('Type', '')

            # Check if this is an image relationship
            if 'image' in rel_type.lower() or (target and '../media/' in target):
                # Extract just the filename from the target path
                media_name = Path(target).name
                rid_to_media[rel_id] = media_name

        # Parse the slide XML to find which relationship IDs are used
        slide_content = zip_ref.read(slide_file).decode('utf-8')
        slide_root = ET.fromstring(slide_content)

        # Find all elements with r:embed or r:link attributes (image references)
        for elem in slide_root.iter():
            for attr_name in ['embed', 'link']:
                # Check with namespace prefix
                rid = elem.get(f'{{{NAMESPACES["r"]}}}{attr_name}')
                if rid and rid in rid_to_media:
                    media_name = rid_to_media[rid]
                    if media_name not in media_to_slides:
                        media_to_slides[media_name] = []
                    if slide_num not in media_to_slides[media_name]:
                        media_to_slides[media_name].append(slide_num)

    # Sort slide numbers for each media file
    for media_name in media_to_slides:
        media_to_slides[media_name].sort()

    return media_to_slides

def extract_pptx_images(pptx_path: Path, output_base: Path):
    """Extract all images from a PPTX file with slide number tracking"""

    relative_path = pptx_path.relative_to(SOURCE_DIR)
    safe_name = str(relative_path).replace("\\", "_").replace("/", "_").replace(" ", "_")
    safe_name = safe_name.replace(".pptx", "")

    output_folder = output_base / safe_name
    output_folder.mkdir(parents=True, exist_ok=True)

    extracted_images = []

    try:
        with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
            # First, build the slide-to-image mapping
            media_to_slides = build_slide_image_mapping(zip_ref)

            # Get all media files
            media_files = [f for f in zip_ref.namelist() if f.startswith('ppt/media/')]

            for i, media_file in enumerate(media_files):
                ext = Path(media_file).suffix.lower()
                if ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.emf', '.wmf']:
                    # Extract image
                    data = zip_ref.read(media_file)

                    # Get the media filename
                    media_name = Path(media_file).name

                    # Look up which slides this image appears on
                    slide_numbers = media_to_slides.get(media_name, [])

                    # Create output filename with slide info
                    output_name = f"image_{i+1:03d}{ext}"
                    output_path = output_folder / output_name

                    with open(output_path, 'wb') as f:
                        f.write(data)

                    extracted_images.append({
                        "source_file": str(pptx_path),
                        "image_path": str(output_path),
                        "original_name": media_name,
                        "index": i + 1,
                        "slide_numbers": slide_numbers,
                        "primary_slide": slide_numbers[0] if slide_numbers else None
                    })

    except Exception as e:
        print(f"  Error extracting from {pptx_path.name}: {e}")
        import traceback
        traceback.print_exc()

    return extracted_images

def extract_slide_text(pptx_path: Path):
    """Extract text content from each slide"""
    slides_text = []

    try:
        prs = Presentation(pptx_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            slides_text.append({
                "slide_number": slide_num,
                "text_content": "\n".join(slide_text)
            })
    except Exception as e:
        print(f"  Error reading slides from {pptx_path.name}: {e}")

    return slides_text

def extract_video_keyframes(video_path: Path, output_base: Path):
    """Extract keyframes from video using ffmpeg"""

    relative_path = video_path.relative_to(SOURCE_DIR)
    safe_name = str(relative_path).replace("\\", "_").replace("/", "_").replace(" ", "_")
    safe_name = Path(safe_name).stem

    output_folder = output_base / safe_name
    output_folder.mkdir(parents=True, exist_ok=True)

    keyframes = []

    try:
        # Get video duration
        probe_cmd = [
            'ffprobe', '-v', 'error', '-show_entries', 'format=duration',
            '-of', 'default=noprint_wrappers=1:nokey=1', str(video_path)
        ]
        result = subprocess.run(probe_cmd, capture_output=True, text=True)
        duration = float(result.stdout.strip())

        # Extract keyframes at intervals
        output_pattern = output_folder / "keyframe_%03d.jpg"

        extract_cmd = [
            'ffmpeg', '-i', str(video_path),
            '-vf', f'fps=1/{KEYFRAME_INTERVAL}',  # 1 frame every N seconds
            '-q:v', '2',  # High quality JPEG
            '-y',  # Overwrite
            str(output_pattern)
        ]

        subprocess.run(extract_cmd, capture_output=True)

        # List extracted keyframes
        for kf in sorted(output_folder.glob("keyframe_*.jpg")):
            frame_num = int(kf.stem.split("_")[1])
            timestamp = (frame_num - 1) * KEYFRAME_INTERVAL

            keyframes.append({
                "source_file": str(video_path),
                "keyframe_path": str(kf),
                "timestamp_seconds": timestamp,
                "timestamp_formatted": f"{timestamp // 60}:{timestamp % 60:02d}"
            })

        print(f"  Extracted {len(keyframes)} keyframes from {video_path.name}")

    except Exception as e:
        print(f"  Error extracting keyframes from {video_path.name}: {e}")

    return keyframes, duration if 'duration' in dir() else 0

def extract_video_audio(video_path: Path, output_base: Path):
    """Extract audio track from video for transcription"""

    safe_name = video_path.stem.replace(" ", "_")
    audio_path = output_base / f"{safe_name}.wav"

    try:
        cmd = [
            'ffmpeg', '-i', str(video_path),
            '-vn',  # No video
            '-acodec', 'pcm_s16le',  # WAV format
            '-ar', '16000',  # 16kHz for Whisper
            '-ac', '1',  # Mono
            '-y',
            str(audio_path)
        ]
        subprocess.run(cmd, capture_output=True)
        print(f"  Extracted audio: {audio_path.name}")
        return audio_path
    except Exception as e:
        print(f"  Error extracting audio from {video_path.name}: {e}")
        return None

def main():
    print("=" * 60)
    print("MEDIA EXTRACTION PIPELINE")
    print("=" * 60)

    # Setup
    pptx_img_dir, keyframe_dir, transcript_dir, metadata_dir = setup_output_dirs()

    all_metadata = {
        "pptx_files": [],
        "video_files": [],
        "total_images_extracted": 0,
        "total_keyframes_extracted": 0
    }

    # Find all PPTX files
    pptx_files = list(SOURCE_DIR.rglob("*.pptx"))
    print(f"\nFound {len(pptx_files)} PPTX files")

    # Process PPTX files
    print("\n" + "-" * 40)
    print("EXTRACTING PPTX IMAGES")
    print("-" * 40)

    total_with_slides = 0
    total_without_slides = 0

    for pptx_path in pptx_files:
        print(f"\nProcessing: {pptx_path.name}")

        # Extract images with slide tracking
        images = extract_pptx_images(pptx_path, pptx_img_dir)

        # Count images with/without slide mapping
        with_slides = sum(1 for img in images if img.get("slide_numbers"))
        without_slides = len(images) - with_slides
        total_with_slides += with_slides
        total_without_slides += without_slides

        print(f"  Extracted {len(images)} images ({with_slides} with slide numbers, {without_slides} without)")

        # Extract slide text
        slides = extract_slide_text(pptx_path)

        all_metadata["pptx_files"].append({
            "source_path": str(pptx_path),
            "relative_path": str(pptx_path.relative_to(SOURCE_DIR)),
            "images_extracted": len(images),
            "images_with_slide_mapping": with_slides,
            "slides_count": len(slides),
            "images": images,
            "slides": slides
        })
        all_metadata["total_images_extracted"] += len(images)

    # Find all video files
    video_extensions = ['.mp4', '.mov', '.avi', '.wmv', '.webm', '.mkv']
    video_files = [f for f in SOURCE_DIR.rglob("*") if f.suffix.lower() in video_extensions]
    print(f"\n\nFound {len(video_files)} video files")

    # Process video files
    print("\n" + "-" * 40)
    print("EXTRACTING VIDEO KEYFRAMES & AUDIO")
    print("-" * 40)

    for video_path in video_files:
        print(f"\nProcessing: {video_path.name}")

        # Extract keyframes
        keyframes, duration = extract_video_keyframes(video_path, keyframe_dir)

        # Extract audio for transcription
        audio_path = extract_video_audio(video_path, transcript_dir)

        all_metadata["video_files"].append({
            "source_path": str(video_path),
            "relative_path": str(video_path.relative_to(SOURCE_DIR)),
            "duration_seconds": duration,
            "keyframes_extracted": len(keyframes),
            "keyframes": keyframes,
            "audio_path": str(audio_path) if audio_path else None
        })
        all_metadata["total_keyframes_extracted"] += len(keyframes)

    # Save metadata
    metadata_path = metadata_dir / "extraction_metadata.json"
    with open(metadata_path, 'w', encoding='utf-8') as f:
        json.dump(all_metadata, f, indent=2, ensure_ascii=False)

    # Add slide mapping stats to metadata
    all_metadata["images_with_slide_mapping"] = total_with_slides
    all_metadata["images_without_slide_mapping"] = total_without_slides

    # Summary
    print("\n" + "=" * 60)
    print("EXTRACTION COMPLETE")
    print("=" * 60)
    print(f"PPTX files processed: {len(all_metadata['pptx_files'])}")
    print(f"Total images extracted: {all_metadata['total_images_extracted']}")
    print(f"  - With slide numbers: {total_with_slides}")
    print(f"  - Without slide numbers: {total_without_slides}")
    print(f"Video files processed: {len(all_metadata['video_files'])}")
    print(f"Total keyframes extracted: {all_metadata['total_keyframes_extracted']}")
    print(f"\nOutput directory: {OUTPUT_DIR}")
    print(f"Metadata saved to: {metadata_path}")

if __name__ == "__main__":
    main()
