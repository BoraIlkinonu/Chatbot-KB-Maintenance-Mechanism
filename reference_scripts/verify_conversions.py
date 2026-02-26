"""
Verification Script
Compares original files with converted versions to check content preservation
"""

import os
from pathlib import Path
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader
from openpyxl import load_workbook
import re

BASE_DIR = Path(r"D:\Term 3 QA\Teacher Resources - Term 2\Teacher Resources")
OUTPUT_DIR = Path(r"D:\Term 3 QA\Teacher Resources - Term 2\Converted")

def normalize_text(text):
    """Normalize text for comparison - remove extra whitespace, lowercase"""
    text = re.sub(r'\s+', ' ', str(text).lower().strip())
    # Remove common markdown artifacts
    text = re.sub(r'[#\-|]', '', text)
    return text

def get_docx_text(docx_path):
    """Extract all text from DOCX"""
    doc = Document(docx_path)
    text_parts = []
    for para in doc.paragraphs:
        text_parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text_parts.append(cell.text)
    return normalize_text(' '.join(text_parts))

def get_pptx_text(pptx_path):
    """Extract all text from PPTX"""
    prs = Presentation(pptx_path)
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_parts.append(shape.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text_parts.append(cell.text)
    return normalize_text(' '.join(text_parts))

def get_pdf_text(pdf_path):
    """Extract all text from PDF"""
    reader = PdfReader(pdf_path)
    text_parts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            text_parts.append(text)
    return normalize_text(' '.join(text_parts))

def get_xlsx_text(xlsx_path):
    """Extract all text from XLSX"""
    wb = load_workbook(xlsx_path, data_only=True)
    text_parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    text_parts.append(str(cell))
    return normalize_text(' '.join(text_parts))

def get_md_text(md_path):
    """Read text from markdown file"""
    with open(md_path, 'r', encoding='utf-8') as f:
        return normalize_text(f.read())

def get_csv_text(csv_dir, base_name):
    """Read text from CSV file(s)"""
    text_parts = []
    for csv_file in csv_dir.glob(f"{base_name}*.csv"):
        with open(csv_file, 'r', encoding='utf-8') as f:
            text_parts.append(f.read())
    return normalize_text(' '.join(text_parts))

def calculate_similarity(original, converted):
    """Calculate text overlap percentage"""
    if not original:
        return 100.0 if not converted else 0.0

    # Split into words
    orig_words = set(original.split())
    conv_words = set(converted.split())

    if not orig_words:
        return 100.0

    # Calculate overlap
    common = orig_words.intersection(conv_words)
    similarity = (len(common) / len(orig_words)) * 100
    return similarity

def verify_all():
    """Verify all conversions"""
    results = []

    for file_path in BASE_DIR.rglob("*"):
        if not file_path.is_file():
            continue
        if file_path.name.startswith("~$"):
            continue

        ext = file_path.suffix.lower()
        rel_path = file_path.relative_to(BASE_DIR)

        try:
            if ext == ".docx":
                original_text = get_docx_text(file_path)
                converted_path = OUTPUT_DIR / rel_path.parent / (file_path.stem + ".md")
                if converted_path.exists():
                    converted_text = get_md_text(converted_path)
                    similarity = calculate_similarity(original_text, converted_text)
                    results.append(("DOCX", file_path.name, similarity, len(original_text), len(converted_text)))

            elif ext == ".pptx":
                original_text = get_pptx_text(file_path)
                converted_path = OUTPUT_DIR / rel_path.parent / (file_path.stem + ".md")
                if converted_path.exists():
                    converted_text = get_md_text(converted_path)
                    similarity = calculate_similarity(original_text, converted_text)
                    results.append(("PPTX", file_path.name, similarity, len(original_text), len(converted_text)))

            elif ext == ".pdf":
                original_text = get_pdf_text(file_path)
                converted_path = OUTPUT_DIR / rel_path.parent / (file_path.stem + ".md")
                if converted_path.exists():
                    converted_text = get_md_text(converted_path)
                    similarity = calculate_similarity(original_text, converted_text)
                    results.append(("PDF", file_path.name, similarity, len(original_text), len(converted_text)))

            elif ext == ".xlsx":
                original_text = get_xlsx_text(file_path)
                converted_dir = OUTPUT_DIR / rel_path.parent
                converted_text = get_csv_text(converted_dir, file_path.stem)
                similarity = calculate_similarity(original_text, converted_text)
                results.append(("XLSX", file_path.name, similarity, len(original_text), len(converted_text)))

        except Exception as e:
            results.append((ext.upper()[1:], file_path.name, -1, 0, 0, str(e)))

    # Print results
    print("="*80)
    print("VERIFICATION RESULTS")
    print("="*80)
    print(f"{'Type':<6} {'File':<50} {'Match%':<8} {'Status'}")
    print("-"*80)

    passed = 0
    warnings = 0
    failed = 0

    for result in sorted(results, key=lambda x: x[2]):
        file_type, name, similarity, orig_len, conv_len = result[:5]
        short_name = name[:47] + "..." if len(name) > 50 else name

        if similarity < 0:
            status = "ERROR"
            failed += 1
        elif similarity >= 90:
            status = "PASS"
            passed += 1
        elif similarity >= 70:
            status = "WARN"
            warnings += 1
        else:
            status = "LOW"
            warnings += 1

        print(f"{file_type:<6} {short_name:<50} {similarity:>6.1f}%  {status}")

    print("-"*80)
    print(f"Summary: {passed} passed, {warnings} warnings, {failed} errors")
    print(f"Total files verified: {len(results)}")

    # Show detailed issues for low matches
    low_matches = [r for r in results if 0 <= r[2] < 80]
    if low_matches:
        print("\n" + "="*80)
        print("FILES WITH LOW CONTENT MATCH (< 80%)")
        print("="*80)
        for r in low_matches:
            print(f"\n{r[1]}")
            print(f"  Original chars: {r[3]}, Converted chars: {r[4]}")
            print(f"  Match: {r[2]:.1f}%")

if __name__ == "__main__":
    verify_all()
