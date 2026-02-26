"""
Document Conversion Script
Converts DOCX, PPTX, PDF, and XLSX files to text-friendly formats (Markdown/CSV)
"""

import os
from pathlib import Path
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader
from openpyxl import load_workbook
import csv
import traceback

BASE_DIR = Path(r"D:\Term 3 QA\Teacher Resources - Term 2\Teacher Resources")
OUTPUT_DIR = Path(r"D:\Term 3 QA\Teacher Resources - Term 2\Converted")

def ensure_output_dir(original_path):
    """Create corresponding output directory structure"""
    rel_path = original_path.relative_to(BASE_DIR)
    output_path = OUTPUT_DIR / rel_path.parent
    output_path.mkdir(parents=True, exist_ok=True)
    return output_path

def convert_docx_to_md(docx_path):
    """Convert DOCX to Markdown"""
    try:
        doc = Document(docx_path)
        output_dir = ensure_output_dir(docx_path)
        output_file = output_dir / (docx_path.stem + ".md")

        content = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Check for headings based on style
                style_name = para.style.name.lower() if para.style else ""
                if "heading 1" in style_name:
                    content.append(f"# {text}\n")
                elif "heading 2" in style_name:
                    content.append(f"## {text}\n")
                elif "heading 3" in style_name:
                    content.append(f"### {text}\n")
                elif "title" in style_name:
                    content.append(f"# {text}\n")
                else:
                    content.append(f"{text}\n")

        # Also extract tables
        for table in doc.tables:
            content.append("\n| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |\n")
            for row in table.rows:
                row_text = " | ".join(cell.text.strip().replace("\n", " ") for cell in row.cells)
                content.append(f"| {row_text} |\n")
            content.append("\n")

        with open(output_file, 'w', encoding='utf-8') as f:
            f.write("\n".join(content))

        return str(output_file), True, None
    except Exception as e:
        return str(docx_path), False, str(e)

def convert_pptx_to_md(pptx_path):
    """Convert PPTX to Markdown"""
    try:
        prs = Presentation(pptx_path)
        output_dir = ensure_output_dir(pptx_path)
        output_file = output_dir / (pptx_path.stem + ".md")

        content = []
        content.append(f"# {pptx_path.stem}\n\n")

        for slide_num, slide in enumerate(prs.slides, 1):
            content.append(f"## Slide {slide_num}\n\n")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    # First text on slide often is title
                    content.append(f"{text}\n\n")

                # Handle tables in slides
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip().replace("\n", " ") for cell in row.cells)
                        content.append(f"| {row_text} |\n")
                    content.append("\n")

            content.append("---\n\n")

        with open(output_file, 'w', encoding='utf-8') as f:
            f.write("".join(content))

        return str(output_file), True, None
    except Exception as e:
        return str(pptx_path), False, str(e)

def convert_pdf_to_md(pdf_path):
    """Convert PDF to Markdown"""
    try:
        reader = PdfReader(pdf_path)
        output_dir = ensure_output_dir(pdf_path)
        output_file = output_dir / (pdf_path.stem + ".md")

        content = []
        content.append(f"# {pdf_path.stem}\n\n")

        for page_num, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text:
                content.append(f"## Page {page_num}\n\n")
                content.append(text + "\n\n")
                content.append("---\n\n")

        with open(output_file, 'w', encoding='utf-8') as f:
            f.write("".join(content))

        return str(output_file), True, None
    except Exception as e:
        return str(pdf_path), False, str(e)

def convert_xlsx_to_csv(xlsx_path):
    """Convert XLSX to CSV"""
    try:
        wb = load_workbook(xlsx_path, data_only=True)
        output_dir = ensure_output_dir(xlsx_path)

        output_files = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            # Clean sheet name for filename
            clean_name = "".join(c if c.isalnum() or c in " -_" else "_" for c in sheet_name)
            output_file = output_dir / f"{xlsx_path.stem}_{clean_name}.csv"

            with open(output_file, 'w', newline='', encoding='utf-8') as f:
                writer = csv.writer(f)
                for row in sheet.iter_rows(values_only=True):
                    # Filter out completely empty rows
                    if any(cell is not None for cell in row):
                        writer.writerow([str(cell) if cell is not None else "" for cell in row])

            output_files.append(str(output_file))

        return output_files, True, None
    except Exception as e:
        return str(xlsx_path), False, str(e)

def main():
    """Main conversion function"""
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    results = {
        "docx": {"success": [], "failed": []},
        "pptx": {"success": [], "failed": []},
        "pdf": {"success": [], "failed": []},
        "xlsx": {"success": [], "failed": []}
    }

    # Find all files
    all_files = list(BASE_DIR.rglob("*"))

    for file_path in all_files:
        if not file_path.is_file():
            continue

        ext = file_path.suffix.lower()

        if ext == ".docx" and not file_path.name.startswith("~$"):
            print(f"Converting DOCX: {file_path.name}")
            output, success, error = convert_docx_to_md(file_path)
            if success:
                results["docx"]["success"].append((str(file_path), output))
            else:
                results["docx"]["failed"].append((str(file_path), error))

        elif ext == ".pptx" and not file_path.name.startswith("~$"):
            print(f"Converting PPTX: {file_path.name}")
            output, success, error = convert_pptx_to_md(file_path)
            if success:
                results["pptx"]["success"].append((str(file_path), output))
            else:
                results["pptx"]["failed"].append((str(file_path), error))

        elif ext == ".pdf":
            print(f"Converting PDF: {file_path.name}")
            output, success, error = convert_pdf_to_md(file_path)
            if success:
                results["pdf"]["success"].append((str(file_path), output))
            else:
                results["pdf"]["failed"].append((str(file_path), error))

        elif ext == ".xlsx" and not file_path.name.startswith("~$"):
            print(f"Converting XLSX: {file_path.name}")
            output, success, error = convert_xlsx_to_csv(file_path)
            if success:
                results["xlsx"]["success"].append((str(file_path), output))
            else:
                results["xlsx"]["failed"].append((str(file_path), error))

    # Print summary
    print("\n" + "="*60)
    print("CONVERSION SUMMARY")
    print("="*60)

    for file_type, data in results.items():
        print(f"\n{file_type.upper()} Files:")
        print(f"  Successful: {len(data['success'])}")
        print(f"  Failed: {len(data['failed'])}")
        if data['failed']:
            for path, error in data['failed']:
                print(f"    - {path}: {error}")

    total_success = sum(len(d['success']) for d in results.values())
    total_failed = sum(len(d['failed']) for d in results.values())
    print(f"\nTotal: {total_success} successful, {total_failed} failed")

    return results

if __name__ == "__main__":
    results = main()
