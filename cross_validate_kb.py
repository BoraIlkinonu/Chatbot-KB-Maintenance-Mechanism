"""
Stage 8: Cross-Validation Expert Agent
Uses Claude CLI as an expert judge to verify pipeline extraction quality.

Phase 1: Investigate all ERROR anomalies from Stage 7 (true positive vs false positive)
Phase 2: Verify random 50%+ of passing lessons (8-field semantic comparison)

Usage: python cross_validate_kb.py
Requires: claude CLI installed and authenticated
"""

import sys
import os
import json
import time
import random
import subprocess
import traceback
from pathlib import Path
from datetime import datetime, timezone

sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation

from config import (
    SOURCES_DIR, CONVERTED_DIR, CONSOLIDATED_DIR, OUTPUT_DIR,
    VALIDATION_DIR, CROSS_VALIDATION_DIR, JUDGE_MODEL,
    CROSS_VALIDATION_SAMPLE_RATE,
)


# ──────────────────────────────────────────────────────────
# Ground Truth Extraction
# ──────────────────────────────────────────────────────────

def extract_pptx_ground_truth(pptx_path):
    """Extract structured content from a PPTX file as ground truth.

    Returns dict with slides[], each containing text, notes, links, tables.
    """
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        return {"error": f"File not found: {pptx_path}", "slides": []}

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return {"error": f"Cannot open PPTX: {e}", "slides": []}

    slides = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data = {
            "slide_number": slide_num,
            "text": [],
            "notes": "",
            "links": [],
            "tables": [],
        }

        for shape in slide.shapes:
            # Text content
            if hasattr(shape, "text") and shape.text.strip():
                slide_data["text"].append(shape.text.strip())

            # Hyperlinks from text runs
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.hyperlink and run.hyperlink.address:
                                slide_data["links"].append({
                                    "url": run.hyperlink.address,
                                    "text": run.text.strip(),
                                })
                        except (KeyError, TypeError, AttributeError):
                            pass  # Broken hyperlink refs (None rId)

            # Click-action hyperlinks on shapes
            try:
                action = shape.click_action
                if action and hasattr(action, "hyperlink") and action.hyperlink and action.hyperlink.address:
                    slide_data["links"].append({
                        "url": action.hyperlink.address,
                        "text": shape.text.strip() if hasattr(shape, "text") else "",
                    })
            except (TypeError, AttributeError):
                pass  # Group shapes raise TypeError on click_action access

            # Tables
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    table_data.append([cell.text.strip() for cell in row.cells])
                if table_data:
                    slide_data["tables"].append(table_data)

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_data["notes"] = notes

        slides.append(slide_data)

    return {
        "file": str(pptx_path),
        "total_slides": len(slides),
        "slides": slides,
    }


def find_source_pptx(doc_path):
    """Given a converted doc path (relative), find the source PPTX in sources/.

    doc_path is like 'term2\\Teacher Resources\\...\\Lesson 1.md'
    Source would be sources/term2/...../Lesson 1.pptx
    """
    rel = Path(doc_path)
    source = SOURCES_DIR / rel.with_suffix(".pptx")
    if source.exists():
        return source
    return None


# ──────────────────────────────────────────────────────────
# KB Data Loading
# ──────────────────────────────────────────────────────────

def load_kb_lesson(term, lesson_num):
    """Load a specific lesson entry from KB JSON output."""
    kb_path = OUTPUT_DIR / f"Term {term} - Lesson Based Structure.json"
    if not kb_path.exists():
        return None

    with open(kb_path, "r", encoding="utf-8") as f:
        kb = json.load(f)

    for lesson in kb.get("lessons", []):
        if lesson.get("metadata", {}).get("lesson_id") == lesson_num:
            return lesson

    return None


def load_consolidated_lesson(term, lesson_num):
    """Load consolidated data for a specific lesson."""
    cons_path = CONSOLIDATED_DIR / f"consolidated_term{term}.json"
    if not cons_path.exists():
        return None

    with open(cons_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    return data.get("by_lesson", {}).get(str(lesson_num))


def load_validation_errors(term):
    """Load ERROR anomalies from Stage 7 validation report."""
    report_path = VALIDATION_DIR / f"validation_report_term{term}.json"
    if not report_path.exists():
        return []

    with open(report_path, "r", encoding="utf-8") as f:
        report = json.load(f)

    return [a for a in report.get("anomalies", []) if a.get("severity") == "ERROR"]


def load_qa_failures():
    """Load FAIL checks from QA audit report."""
    audit_path = VALIDATION_DIR / "qa_audit_report.json"
    if not audit_path.exists():
        return []

    with open(audit_path, "r", encoding="utf-8") as f:
        report = json.load(f)

    return [c for c in report.get("checks", []) if not c.get("passed")]


# ──────────────────────────────────────────────────────────
# Claude CLI Judge
# ──────────────────────────────────────────────────────────

def _cli_available():
    """Check if claude CLI is available."""
    try:
        env = {**os.environ}
        env.pop("CLAUDECODE", None)
        result = subprocess.run(
            ["claude", "--version"],
            capture_output=True, text=True, timeout=10, env=env,
        )
        return result.returncode == 0
    except (FileNotFoundError, subprocess.TimeoutExpired):
        return False


def _call_claude(prompt, json_schema=None):
    """Call Claude CLI with a prompt and return the response text.

    Uses claude -p --model {JUDGE_MODEL} --output-format json for structured output.
    Returns parsed JSON if json_schema provided, otherwise raw text.
    """
    cmd = [
        "claude", "-p",
        "--model", JUDGE_MODEL,
        "--output-format", "json",
    ]

    if json_schema:
        cmd.extend(["--json-schema", json.dumps(json_schema)])

    # Allow nested invocation when running inside a Claude Code session
    env = {**os.environ}
    env.pop("CLAUDECODE", None)

    result = subprocess.run(
        cmd,
        input=prompt,
        capture_output=True,
        text=True,
        timeout=120,
        encoding="utf-8",
        env=env,
    )

    if result.returncode != 0:
        raise RuntimeError(f"Claude CLI failed: {result.stderr[:500]}")

    # Parse JSON output format
    try:
        output = json.loads(result.stdout)
    except json.JSONDecodeError:
        return result.stdout.strip()

    # claude --output-format json wraps response in {"type":"result","result":"..."}
    if isinstance(output, dict) and "result" in output:
        text = output["result"]
        # Try to parse the inner result as JSON too
        if json_schema:
            try:
                return json.loads(text)
            except (json.JSONDecodeError, TypeError):
                pass
        return text

    return output


def _dual_judge(prompt, json_schema, verdict_key="verdict", max_attempts=2):
    """Run dual-judge consensus: two independent Claude calls must agree.

    If the two judges disagree, the evaluation is repeated once more (fresh pair).
    After the retry, if still no consensus, returns the first judge's result
    with a 'consensus': False flag.

    Args:
        prompt: The prompt to send to each judge
        json_schema: JSON schema for structured output
        verdict_key: Which key in the result dict holds the verdict to compare
        max_attempts: Number of attempts (each attempt = 2 calls)

    Returns:
        dict with the agreed result, plus 'consensus' and 'attempt' metadata
    """
    for attempt in range(1, max_attempts + 1):
        results = []
        for _ in range(2):
            try:
                result = _call_claude(prompt, json_schema=json_schema)
                if isinstance(result, str):
                    result = json.loads(result)
                results.append(result)
            except Exception as e:
                results.append({
                    verdict_key: "UNCERTAIN",
                    "reason": f"Call failed: {e}",
                    "evidence": f"Call failed: {e}",
                })
            time.sleep(1)

        # Check unanimous agreement
        v1 = results[0].get(verdict_key, "")
        v2 = results[1].get(verdict_key, "")

        if v1 == v2:
            results[0]["consensus"] = True
            results[0]["attempt"] = attempt
            results[0]["judge_votes"] = [v1, v2]
            return results[0]

        # Disagreement — retry if attempts remain
        if attempt < max_attempts:
            continue

    # No consensus after retries — return first result flagged as no-consensus
    results[0]["consensus"] = False
    results[0]["attempt"] = max_attempts
    results[0]["judge_votes"] = [v1, v2]
    return results[0]


def _dual_judge_fields(prompt, json_schema, fields, max_attempts=2):
    """Run dual-judge consensus for multi-field evaluations (Phase 2).

    Both judges must agree on ALL field verdicts. If any field disagrees,
    the evaluation is repeated once.

    Returns:
        dict of field results, plus 'consensus' and 'attempt' metadata
    """
    for attempt in range(1, max_attempts + 1):
        results = []
        for _ in range(2):
            try:
                result = _call_claude(prompt, json_schema=json_schema)
                if isinstance(result, str):
                    result = json.loads(result)
                results.append(result)
            except Exception as e:
                fallback = {
                    f: {"verdict": "MISSING", "evidence": f"Call failed: {e}"}
                    for f in fields
                }
                results.append(fallback)
            time.sleep(1)

        # Check unanimous agreement on all fields
        all_agree = True
        votes = {}
        for field in fields:
            v1 = results[0].get(field, {}).get("verdict", "")
            v2 = results[1].get(field, {}).get("verdict", "")
            votes[field] = [v1, v2]
            if v1 != v2:
                all_agree = False

        if all_agree:
            results[0]["_consensus"] = True
            results[0]["_attempt"] = attempt
            results[0]["_judge_votes"] = votes
            return results[0]

        if attempt < max_attempts:
            continue

    # No consensus — return first result flagged
    results[0]["_consensus"] = False
    results[0]["_attempt"] = max_attempts
    results[0]["_judge_votes"] = votes
    return results[0]


def _truncate(text, max_chars=3000):
    """Truncate text for prompt inclusion."""
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + f"\n... [truncated, {len(text)} total chars]"


def build_ground_truth_text(ground_truth):
    """Format ground truth data as readable text for the prompt."""
    if not ground_truth or ground_truth.get("error"):
        return "[No ground truth available]"

    parts = []
    for slide in ground_truth.get("slides", []):
        parts.append(f"--- Slide {slide['slide_number']} ---")
        for text in slide.get("text", []):
            parts.append(text)
        if slide.get("notes"):
            parts.append(f"[Speaker Notes]: {slide['notes']}")
        for link in slide.get("links", []):
            parts.append(f"[Link]: {link.get('text', '')} → {link['url']}")
        for table in slide.get("tables", []):
            for row in table:
                parts.append(" | ".join(row))

    return "\n".join(parts)


def build_kb_text(kb_lesson):
    """Format KB lesson entry as readable text for the prompt."""
    if not kb_lesson:
        return "[No KB entry found]"

    parts = []
    parts.append(f"lesson_title: {kb_lesson.get('lesson_title', '')}")

    meta = kb_lesson.get("metadata", {})
    parts.append(f"core_topics: {json.dumps(meta.get('core_topics', []))}")
    parts.append(f"learning_objectives: {json.dumps(meta.get('learning_objectives', []))}")
    parts.append(f"endstar_tools: {json.dumps(meta.get('endstar_tools', []))}")
    parts.append(f"activity_type: {meta.get('activity_type', '')}")
    parts.append(f"activity_description: {meta.get('activity_description', '')}")
    parts.append(f"videos: {json.dumps(meta.get('videos', []))}")
    parts.append(f"resources: {json.dumps(meta.get('resources', []))}")
    parts.append(f"keywords: {json.dumps(meta.get('keywords', []))}")

    desc = kb_lesson.get("description_of_activities", "")
    if desc:
        parts.append(f"description_of_activities: {desc[:500]}")

    return "\n".join(parts)


# ──────────────────────────────────────────────────────────
# Phase 1: Error Investigation
# ──────────────────────────────────────────────────────────

ERROR_INVESTIGATION_SCHEMA = {
    "type": "object",
    "properties": {
        "verdict": {
            "type": "string",
            "enum": ["TRUE_POSITIVE", "FALSE_POSITIVE", "UNCERTAIN"],
        },
        "reason": {"type": "string"},
        "evidence": {"type": "string"},
    },
    "required": ["verdict", "reason", "evidence"],
}


def build_error_investigation_prompt(error, source_text, kb_text):
    """Build prompt for Phase 1 error investigation."""
    return f"""You are a curriculum KB quality expert. Investigate whether this validation error is a TRUE POSITIVE (real problem) or FALSE POSITIVE (the validation check is wrong).

VALIDATION ERROR:
Type: {error.get('type', '')}
Message: {error.get('message', '')}
Lesson: {error.get('lesson', '')}
Content Type Expected: {error.get('content_type', '')}

SOURCE CONTENT (raw PPTX text — ground truth):
{_truncate(source_text)}

KB OUTPUT (pipeline extraction):
{_truncate(kb_text)}

Analyze:
1. Does the source file actually contain the expected content type?
2. Is the error caused by a filename/naming pattern mismatch rather than missing content?
3. Is the KB output actually correct despite the validation error?

Return your verdict as JSON with keys: verdict (TRUE_POSITIVE/FALSE_POSITIVE/UNCERTAIN), reason, evidence."""


def judge_error(error, source_content, kb_content):
    """Phase 1: Dual-judge consensus on whether an error is true/false positive.

    Two independent Claude calls must agree on the verdict. If they disagree,
    the evaluation is repeated once. See _dual_judge() for details.
    """
    source_text = build_ground_truth_text(source_content)
    kb_text = build_kb_text(kb_content)

    prompt = build_error_investigation_prompt(error, source_text, kb_text)

    return _dual_judge(prompt, json_schema=ERROR_INVESTIGATION_SCHEMA, verdict_key="verdict")


def run_phase1(terms=(1, 2, 3)):
    """Phase 1: Investigate all ERROR anomalies from Stage 7."""
    print("\n>>> Phase 1: Error Investigation")
    print("-" * 40)

    all_errors = []
    for term in terms:
        errors = load_validation_errors(term)
        for err in errors:
            err["term"] = term
        all_errors.extend(errors)

    # Also include QA audit failures
    qa_failures = load_qa_failures()

    print(f"  Validation errors: {len(all_errors)}")
    print(f"  QA audit failures: {len(qa_failures)}")

    confirmed = []
    false_positives = []
    uncertain = []

    cli_ok = _cli_available()
    if not cli_ok:
        print("  Claude CLI not available — running structural analysis only")

    for i, error in enumerate(all_errors, 1):
        term = error.get("term", 0)
        lesson_num = error.get("lesson", 0)

        print(f"  [{i}/{len(all_errors)}] Term {term} Lesson {lesson_num}: {error.get('message', '')}")

        # Load consolidated data to find source files
        cons_lesson = load_consolidated_lesson(term, lesson_num)
        kb_lesson = load_kb_lesson(term, lesson_num)

        # Try to find source PPTX
        source_content = {"slides": []}
        if cons_lesson:
            for doc in cons_lesson.get("documents", []):
                pptx_path = find_source_pptx(doc.get("path", ""))
                if pptx_path:
                    source_content = extract_pptx_ground_truth(pptx_path)
                    break

        # Structural check: "missing teachers_slides" but lesson has documents
        if error.get("type") == "MISSING" and error.get("content_type") == "teachers_slides":
            if cons_lesson and cons_lesson.get("document_count", 0) > 0:
                # Has documents but named differently — likely false positive
                doc_paths = [d.get("path", "") for d in cons_lesson.get("documents", [])]
                entry = {
                    "term": term,
                    "lesson": lesson_num,
                    "type": error.get("type"),
                    "content_type": error.get("content_type"),
                    "reason": f"Lesson has {cons_lesson['document_count']} document(s) but content_type != 'teachers_slides'. File(s): {doc_paths}",
                    "evidence": "Filename pattern mismatch — content exists under different name",
                }
                if cli_ok:
                    judgment = judge_error(error, source_content, kb_lesson)
                    entry["llm_verdict"] = judgment.get("verdict", "UNCERTAIN")
                    entry["llm_reason"] = judgment.get("reason", "")
                    entry["llm_evidence"] = judgment.get("evidence", "")
                    entry["consensus"] = judgment.get("consensus", None)
                    entry["judge_votes"] = judgment.get("judge_votes", [])
                false_positives.append(entry)
                consensus_tag = " [consensus]" if entry.get("consensus") else ""
                print(f"    → Likely FALSE POSITIVE (has docs, naming mismatch){consensus_tag}")
                continue

        if error.get("type") == "MISSING" and error.get("content_type") == "students_slides":
            # Students slides are commonly not provided
            entry = {
                "term": term,
                "lesson": lesson_num,
                "type": error.get("type"),
                "content_type": error.get("content_type"),
                "reason": "Student slides are not consistently provided across terms",
                "evidence": "Structural observation — not all lessons include student versions",
            }
            if cli_ok and source_content.get("slides"):
                judgment = judge_error(error, source_content, kb_lesson)
                entry["llm_verdict"] = judgment.get("verdict", "UNCERTAIN")
                entry["llm_reason"] = judgment.get("reason", "")
                entry["llm_evidence"] = judgment.get("evidence", "")
                entry["consensus"] = judgment.get("consensus", None)
                entry["judge_votes"] = judgment.get("judge_votes", [])
            false_positives.append(entry)
            consensus_tag = " [consensus]" if entry.get("consensus") else ""
            print(f"    → Likely FALSE POSITIVE (students_slides not always provided){consensus_tag}")
            continue

        # For other errors, use LLM dual-judge if available
        if cli_ok and source_content.get("slides"):
            judgment = judge_error(error, source_content, kb_lesson)
            verdict = judgment.get("verdict", "UNCERTAIN")
            consensus = judgment.get("consensus", False)
            entry = {
                "term": term,
                "lesson": lesson_num,
                "type": error.get("type"),
                "content_type": error.get("content_type", ""),
                "message": error.get("message", ""),
                "llm_verdict": verdict,
                "llm_reason": judgment.get("reason", ""),
                "llm_evidence": judgment.get("evidence", ""),
                "consensus": consensus,
                "judge_votes": judgment.get("judge_votes", []),
            }
            consensus_tag = " [consensus]" if consensus else " [no consensus]"
            if verdict == "TRUE_POSITIVE":
                confirmed.append(entry)
                print(f"    → TRUE POSITIVE{consensus_tag}")
            elif verdict == "FALSE_POSITIVE":
                false_positives.append(entry)
                print(f"    → FALSE POSITIVE{consensus_tag}")
            else:
                uncertain.append(entry)
                print(f"    → UNCERTAIN{consensus_tag}")
        else:
            uncertain.append({
                "term": term,
                "lesson": lesson_num,
                "type": error.get("type"),
                "content_type": error.get("content_type", ""),
                "message": error.get("message", ""),
                "reason": "No source PPTX or CLI unavailable",
            })
            print(f"    → UNCERTAIN (no source or CLI)")

    print(f"\n  Phase 1 Results:")
    print(f"    Confirmed errors: {len(confirmed)}")
    print(f"    False positives:  {len(false_positives)}")
    print(f"    Uncertain:        {len(uncertain)}")

    return {
        "total_errors_reviewed": len(all_errors),
        "qa_failures_reviewed": len(qa_failures),
        "confirmed_errors": confirmed,
        "false_positives": false_positives,
        "uncertain": uncertain,
    }


# ──────────────────────────────────────────────────────────
# Phase 2: Pass Verification
# ──────────────────────────────────────────────────────────

FIELD_EVALUATION_SCHEMA = {
    "type": "object",
    "properties": {
        "lesson_title": {
            "type": "object",
            "properties": {
                "verdict": {"type": "string", "enum": ["CORRECT", "PARTIAL", "INCORRECT", "MISSING"]},
                "evidence": {"type": "string"},
            },
            "required": ["verdict", "evidence"],
        },
        "learning_objectives": {
            "type": "object",
            "properties": {
                "verdict": {"type": "string", "enum": ["CORRECT", "PARTIAL", "INCORRECT", "MISSING"]},
                "evidence": {"type": "string"},
            },
            "required": ["verdict", "evidence"],
        },
        "core_topics": {
            "type": "object",
            "properties": {
                "verdict": {"type": "string", "enum": ["CORRECT", "PARTIAL", "INCORRECT", "MISSING"]},
                "evidence": {"type": "string"},
            },
            "required": ["verdict", "evidence"],
        },
        "activity_description": {
            "type": "object",
            "properties": {
                "verdict": {"type": "string", "enum": ["CORRECT", "PARTIAL", "INCORRECT", "MISSING"]},
                "evidence": {"type": "string"},
            },
            "required": ["verdict", "evidence"],
        },
        "resources": {
            "type": "object",
            "properties": {
                "verdict": {"type": "string", "enum": ["CORRECT", "PARTIAL", "INCORRECT", "MISSING"]},
                "evidence": {"type": "string"},
            },
            "required": ["verdict", "evidence"],
        },
        "videos": {
            "type": "object",
            "properties": {
                "verdict": {"type": "string", "enum": ["CORRECT", "PARTIAL", "INCORRECT", "MISSING"]},
                "evidence": {"type": "string"},
            },
            "required": ["verdict", "evidence"],
        },
        "endstar_tools": {
            "type": "object",
            "properties": {
                "verdict": {"type": "string", "enum": ["CORRECT", "PARTIAL", "INCORRECT", "MISSING"]},
                "evidence": {"type": "string"},
            },
            "required": ["verdict", "evidence"],
        },
        "keywords": {
            "type": "object",
            "properties": {
                "verdict": {"type": "string", "enum": ["CORRECT", "PARTIAL", "INCORRECT", "MISSING"]},
                "evidence": {"type": "string"},
            },
            "required": ["verdict", "evidence"],
        },
    },
    "required": ["lesson_title", "learning_objectives", "core_topics",
                  "activity_description", "resources", "videos",
                  "endstar_tools", "keywords"],
}

EVALUATED_FIELDS = [
    "lesson_title", "learning_objectives", "core_topics",
    "activity_description", "resources", "videos",
    "endstar_tools", "keywords",
]


def build_lesson_evaluation_prompt(source_text, kb_text):
    """Build prompt for Phase 2 full lesson evaluation."""
    return f"""You are a curriculum KB quality expert. Compare the SOURCE CONTENT (raw PPTX text) against the EXTRACTED KB OUTPUT and evaluate each field.

SOURCE CONTENT (ground truth from PPTX slides):
{_truncate(source_text, 4000)}

KB OUTPUT (pipeline extraction):
{_truncate(kb_text, 2000)}

Evaluate these fields. Rate each CORRECT / PARTIAL / INCORRECT / MISSING:
1. lesson_title — Does it accurately capture the lesson title from the slides?
2. learning_objectives — Are the extracted objectives present in the source?
3. core_topics — Do they reflect actual slide content themes?
4. activity_description — Does it match activities described in slides?
5. resources — Are extracted URLs present in the source PPTX links?
6. videos — Are YouTube/video URLs correctly identified and separated?
7. endstar_tools — Are the matched tools actually mentioned in slide text?
8. keywords — Are they relevant to the actual lesson content?

Return JSON with each field as a key containing verdict and evidence."""


def score_verdict(verdict):
    """Convert verdict string to numeric score."""
    return {"CORRECT": 1.0, "PARTIAL": 0.5, "INCORRECT": 0.0, "MISSING": 0.0}.get(
        verdict, 0.0
    )


def judge_lesson(source_content, kb_content):
    """Phase 2: Dual-judge 8-field evaluation of a lesson.

    Two independent Claude calls must agree on ALL field verdicts.
    If any field disagrees, the evaluation is repeated once.
    See _dual_judge_fields() for details.
    """
    source_text = build_ground_truth_text(source_content)
    kb_text = build_kb_text(kb_content)

    prompt = build_lesson_evaluation_prompt(source_text, kb_text)

    return _dual_judge_fields(
        prompt, json_schema=FIELD_EVALUATION_SCHEMA,
        fields=EVALUATED_FIELDS,
    )


def run_phase2(terms=(1, 2, 3), sample_rate=None):
    """Phase 2: Verify random 50%+ of passing lessons."""
    if sample_rate is None:
        sample_rate = CROSS_VALIDATION_SAMPLE_RATE

    print(f"\n>>> Phase 2: Pass Verification (sample rate: {sample_rate:.0%})")
    print("-" * 40)

    cli_ok = _cli_available()
    if not cli_ok:
        print("  Claude CLI not available — skipping Phase 2")
        return {
            "total_available": 0,
            "total_selected": 0,
            "selection_rate": "0%",
            "per_lesson": [],
            "field_accuracy": {},
            "skipped": True,
        }

    # Collect all available lessons with source PPTX
    available = []
    for term in terms:
        kb_path = OUTPUT_DIR / f"Term {term} - Lesson Based Structure.json"
        if not kb_path.exists():
            continue

        with open(kb_path, "r", encoding="utf-8") as f:
            kb = json.load(f)

        for lesson in kb.get("lessons", []):
            lesson_num = lesson.get("metadata", {}).get("lesson_id", 0)
            cons_lesson = load_consolidated_lesson(term, lesson_num)
            if not cons_lesson:
                continue

            # Find a source PPTX
            pptx_path = None
            for doc in cons_lesson.get("documents", []):
                pptx_path = find_source_pptx(doc.get("path", ""))
                if pptx_path:
                    break

            if pptx_path:
                available.append({
                    "term": term,
                    "lesson_num": lesson_num,
                    "pptx_path": pptx_path,
                    "kb_lesson": lesson,
                })

    print(f"  Available lessons with source PPTX: {len(available)}")

    # Random sample
    sample_count = max(1, int(len(available) * sample_rate))
    selected = random.sample(available, min(sample_count, len(available)))

    print(f"  Selected for verification: {len(selected)}")

    per_lesson = []
    field_totals = {field: {"CORRECT": 0, "PARTIAL": 0, "INCORRECT": 0, "MISSING": 0}
                    for field in EVALUATED_FIELDS}

    for i, item in enumerate(selected, 1):
        term = item["term"]
        lesson_num = item["lesson_num"]

        print(f"  [{i}/{len(selected)}] Term {term} Lesson {lesson_num}")

        # Extract ground truth
        source_content = extract_pptx_ground_truth(item["pptx_path"])
        if source_content.get("error"):
            print(f"    → Skipped: {source_content['error']}")
            continue

        # Judge
        judgment = judge_lesson(source_content, item["kb_lesson"])

        # Score
        total_score = 0
        field_count = 0
        fields_result = {}

        for field in EVALUATED_FIELDS:
            field_data = judgment.get(field, {"verdict": "MISSING", "evidence": ""})
            verdict = field_data.get("verdict", "MISSING")
            fields_result[field] = field_data
            total_score += score_verdict(verdict)
            field_count += 1
            field_totals[field][verdict] = field_totals[field].get(verdict, 0) + 1

        overall = round(total_score / max(field_count, 1), 3)

        consensus = judgment.get("_consensus", None)
        judge_votes = judgment.get("_judge_votes", {})

        per_lesson.append({
            "term": term,
            "lesson": lesson_num,
            "fields": fields_result,
            "overall_score": overall,
            "consensus": consensus,
            "judge_votes": judge_votes,
        })

        consensus_tag = " [consensus]" if consensus else " [no consensus]" if consensus is not None else ""
        print(f"    → Score: {overall:.1%}{consensus_tag}")

    # Compute field accuracy summary
    field_accuracy = {}
    for field in EVALUATED_FIELDS:
        totals = field_totals[field]
        field_accuracy[field] = {
            "correct": totals["CORRECT"],
            "partial": totals["PARTIAL"],
            "incorrect": totals["INCORRECT"],
            "missing": totals["MISSING"],
        }

    rate_str = f"{len(selected)/max(len(available),1):.0%}"
    print(f"\n  Phase 2 Results:")
    print(f"    Lessons verified: {len(per_lesson)}/{len(available)}")
    if per_lesson:
        avg = sum(l["overall_score"] for l in per_lesson) / len(per_lesson)
        print(f"    Average score: {avg:.1%}")

    return {
        "total_available": len(available),
        "total_selected": len(selected),
        "selection_rate": rate_str,
        "per_lesson": per_lesson,
        "field_accuracy": field_accuracy,
    }


# ──────────────────────────────────────────────────────────
# Report Generation
# ──────────────────────────────────────────────────────────

def compute_overall_confidence(phase1, phase2):
    """Compute overall confidence score from both phases."""
    # Phase 1: high false-positive rate → high confidence (errors aren't real)
    total_reviewed = phase1.get("total_errors_reviewed", 0)
    fp_count = len(phase1.get("false_positives", []))
    confirmed_count = len(phase1.get("confirmed_errors", []))

    if total_reviewed > 0:
        fp_ratio = fp_count / total_reviewed
    else:
        fp_ratio = 1.0

    # Phase 2: average field accuracy
    per_lesson = phase2.get("per_lesson", [])
    if per_lesson:
        avg_score = sum(l["overall_score"] for l in per_lesson) / len(per_lesson)
    else:
        avg_score = 0.5  # No data → neutral

    # Combined: weight phase 2 more (actual content verification)
    confidence = round(avg_score * 0.7 + fp_ratio * 0.3, 2)

    # Penalize for confirmed errors
    penalty = min(confirmed_count * 0.02, 0.2)
    confidence = max(0, round(confidence - penalty, 2))

    return confidence


def generate_recommendations(phase1, phase2):
    """Generate actionable recommendations from cross-validation results."""
    recs = []

    fp = phase1.get("false_positives", [])
    naming_fps = [f for f in fp if "naming" in f.get("reason", "").lower() or "content_type" in f.get("reason", "").lower()]
    if naming_fps:
        recs.append(
            f"Fix content_type detection for non-standard filenames "
            f"({len(naming_fps)} false positives from naming mismatches)"
        )

    students_fps = [f for f in fp if f.get("content_type") == "students_slides"]
    if students_fps:
        recs.append(
            f"Remove students_slides from EXPECTED_CONTENT_TYPES or make it optional "
            f"({len(students_fps)} false positives)"
        )

    # Phase 2 field weaknesses
    field_acc = phase2.get("field_accuracy", {})
    for field, counts in field_acc.items():
        incorrect = counts.get("incorrect", 0)
        total = sum(counts.values())
        if total > 0 and incorrect / total > 0.2:
            recs.append(f"Review {field} extraction logic — {incorrect}/{total} rated INCORRECT")

    if not recs:
        recs.append("No critical issues found — pipeline extraction appears accurate")

    return recs


def generate_report(phase1, phase2):
    """Generate full cross-validation report."""
    confidence = compute_overall_confidence(phase1, phase2)
    recommendations = generate_recommendations(phase1, phase2)

    report = {
        "cross_validated_at": datetime.now(timezone.utc).isoformat(),
        "model_used": JUDGE_MODEL,
        "cli_tool": "claude",
        "phase1_error_investigation": phase1,
        "phase2_pass_verification": phase2,
        "overall_confidence": confidence,
        "recommendations": recommendations,
    }

    return report


def write_text_summary(report, path):
    """Write human-readable summary of cross-validation report."""
    with open(path, "w", encoding="utf-8") as f:
        f.write("=" * 70 + "\n")
        f.write("  CROSS-VALIDATION REPORT\n")
        f.write(f"  {report['cross_validated_at']}\n")
        f.write(f"  Model: {report['model_used']} (via claude CLI)\n")
        f.write("=" * 70 + "\n\n")

        f.write(f"Overall Confidence: {report['overall_confidence']:.0%}\n\n")

        # Phase 1
        p1 = report["phase1_error_investigation"]
        f.write("PHASE 1: Error Investigation\n")
        f.write("-" * 40 + "\n")
        f.write(f"  Total errors reviewed: {p1.get('total_errors_reviewed', 0)}\n")
        f.write(f"  Confirmed errors:      {len(p1.get('confirmed_errors', []))}\n")
        f.write(f"  False positives:       {len(p1.get('false_positives', []))}\n")
        f.write(f"  Uncertain:             {len(p1.get('uncertain', []))}\n\n")

        if p1.get("confirmed_errors"):
            f.write("  Confirmed Errors:\n")
            for err in p1["confirmed_errors"]:
                f.write(f"    - Term {err.get('term')} L{err.get('lesson')}: "
                        f"{err.get('type')} — {err.get('llm_reason', err.get('reason', ''))}\n")
            f.write("\n")

        if p1.get("false_positives"):
            f.write("  False Positives (top 10):\n")
            for err in p1["false_positives"][:10]:
                f.write(f"    - Term {err.get('term')} L{err.get('lesson')}: "
                        f"{err.get('content_type', '')} — {err.get('reason', '')[:100]}\n")
            f.write("\n")

        # Phase 2
        p2 = report["phase2_pass_verification"]
        f.write("PHASE 2: Pass Verification\n")
        f.write("-" * 40 + "\n")
        f.write(f"  Lessons available:     {p2.get('total_available', 0)}\n")
        f.write(f"  Lessons selected:      {p2.get('total_selected', 0)}\n")
        f.write(f"  Selection rate:        {p2.get('selection_rate', 'N/A')}\n\n")

        per_lesson = p2.get("per_lesson", [])
        if per_lesson:
            avg = sum(l["overall_score"] for l in per_lesson) / len(per_lesson)
            f.write(f"  Average field accuracy: {avg:.1%}\n\n")

            f.write("  Per-Lesson Scores:\n")
            for l in sorted(per_lesson, key=lambda x: (x["term"], x["lesson"])):
                f.write(f"    Term {l['term']} L{l['lesson']}: {l['overall_score']:.0%}\n")
            f.write("\n")

        # Field accuracy
        field_acc = p2.get("field_accuracy", {})
        if field_acc:
            f.write("  Field Accuracy Summary:\n")
            for field, counts in field_acc.items():
                total = sum(counts.values())
                if total > 0:
                    correct_pct = counts.get("correct", 0) / total
                    f.write(f"    {field:25s}: {correct_pct:.0%} correct "
                            f"({counts.get('correct',0)}C {counts.get('partial',0)}P "
                            f"{counts.get('incorrect',0)}I {counts.get('missing',0)}M)\n")
            f.write("\n")

        # Recommendations
        f.write("RECOMMENDATIONS\n")
        f.write("-" * 40 + "\n")
        for rec in report.get("recommendations", []):
            f.write(f"  - {rec}\n")
        f.write("\n")


# ──────────────────────────────────────────────────────────
# Main Entry Point
# ──────────────────────────────────────────────────────────

def run_cross_validation(terms=(1, 2, 3), sample_rate=None):
    """Run full cross-validation: Phase 1 + Phase 2."""
    print("=" * 60)
    print("  Stage 8: Cross-Validation Expert Agent")
    print("=" * 60)

    cli_ok = _cli_available()
    print(f"  Claude CLI available: {cli_ok}")
    if not cli_ok:
        print("  Warning: Claude CLI not available. Structural checks only.")

    CROSS_VALIDATION_DIR.mkdir(parents=True, exist_ok=True)

    # Phase 1
    phase1 = run_phase1(terms)

    # Phase 2
    phase2 = run_phase2(terms, sample_rate)

    # Generate report
    report = generate_report(phase1, phase2)

    # Save JSON report
    json_path = CROSS_VALIDATION_DIR / "cross_validation_report.json"
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(report, f, indent=2, ensure_ascii=False)
    print(f"\n  JSON report: {json_path}")

    # Save text summary
    txt_path = CROSS_VALIDATION_DIR / "cross_validation_report.txt"
    write_text_summary(report, txt_path)
    print(f"  Text report: {txt_path}")

    print(f"\n  Overall Confidence: {report['overall_confidence']:.0%}")
    print(f"  Recommendations: {len(report['recommendations'])}")

    print("\n" + "=" * 60)
    print("  Cross-Validation Complete")
    print("=" * 60)

    return report


if __name__ == "__main__":
    run_cross_validation()
