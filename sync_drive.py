"""
Google Drive Sync Script
Scans target folders, detects changes since last sync, downloads changed files,
logs EVERYTHING (including UNCHANGED), and captures Drive Activity data.
"""

import sys
import os
import json
import hashlib
import io
import time
from datetime import datetime, timezone
from pathlib import Path

sys.stdout.reconfigure(encoding="utf-8")

from googleapiclient.http import MediaIoBaseDownload

from config import (
    TARGET_FOLDERS, SOURCES_DIR, LOGS_DIR, PREVIOUS_SCAN_FILE,
    NATIVE_GOOGLE_MIMES, LOG_EVERYTHING, BASE_DIR, EXPORTS_FOLDER_ID,
)
from auth import authenticate, get_drive_service, get_activity_service


# ──────────────────────────────────────────────────────────
# Drive scanning
# ──────────────────────────────────────────────────────────

def scan_folder(service, folder_id, depth=0, folder_path=""):
    """Recursively scan a Drive folder. Returns list of file metadata dicts.
    Tracks the full folder path so downloaded files preserve hierarchy."""
    files = []
    page_token = None

    while True:
        resp = service.files().list(
            q=f"'{folder_id}' in parents and trashed = false",
            fields=(
                "nextPageToken,files(id,name,mimeType,size,md5Checksum,"
                "createdTime,modifiedTime,lastModifyingUser,owners,"
                "webViewLink,fileExtension,version,headRevisionId,"
                "parents,shared,description)"
            ),
            pageSize=1000,
            pageToken=page_token,
            orderBy="folder,name",
        ).execute()

        for item in resp.get("files", []):
            mime = item.get("mimeType", "")

            if mime == "application/vnd.google-apps.folder":
                # Recurse into subfolder, tracking path
                subfolder_name = item.get("name", "")
                subfolder_path = f"{folder_path}/{subfolder_name}" if folder_path else subfolder_name
                children = scan_folder(service, item["id"], depth + 1, subfolder_path)
                files.extend(children)
            elif mime == "application/vnd.google-apps.shortcut":
                # Shortcuts are pointers to files/folders already in the scan tree.
                # Skip them to avoid duplicate scanning and extra API calls.
                print(f"  Skipping shortcut: {item.get('name', '')}")
            else:
                last_mod = item.get("lastModifyingUser", {})
                owners = item.get("owners", [])

                files.append({
                    "id": item["id"],
                    "name": item.get("name", ""),
                    "mime_type": mime,
                    "size": int(item.get("size", 0) or 0),
                    "md5": item.get("md5Checksum", ""),
                    "created_time": item.get("createdTime", ""),
                    "modified_time": item.get("modifiedTime", ""),
                    "version": item.get("version", ""),
                    "head_revision_id": item.get("headRevisionId", ""),
                    "web_link": item.get("webViewLink", ""),
                    "extension": item.get("fileExtension", ""),
                    "parent_id": folder_id,
                    "folder_path": folder_path,
                    "shared": item.get("shared", False),
                    "description": item.get("description", ""),
                    "last_modifier_email": last_mod.get("emailAddress", ""),
                    "last_modifier_name": last_mod.get("displayName", ""),
                    "owner_email": owners[0].get("emailAddress", "") if owners else "",
                    "owner_name": owners[0].get("displayName", "") if owners else "",
                    "is_native_google": mime.startswith("application/vnd.google-apps."),
                    "native_type": NATIVE_GOOGLE_MIMES.get(mime),
                })

        page_token = resp.get("nextPageToken")
        if not page_token:
            break

    return files


# ──────────────────────────────────────────────────────────
# Change detection
# ──────────────────────────────────────────────────────────

def load_previous_scan():
    """Load previous scan results for comparison."""
    if PREVIOUS_SCAN_FILE.exists():
        with open(PREVIOUS_SCAN_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}


def detect_changes(current_files, previous_files_by_id):
    """
    Compare current scan against previous scan.
    Returns a list of change records for EVERY file (including UNCHANGED).
    """
    changes = []
    current_ids = set()

    for f in current_files:
        fid = f["id"]
        current_ids.add(fid)
        prev = previous_files_by_id.get(fid)

        if prev is None:
            changes.append({**f, "change_type": "NEW"})
        elif f["md5"] and prev.get("md5") and f["md5"] != prev["md5"]:
            changes.append({**f, "change_type": "MODIFIED", "previous_md5": prev["md5"]})
        elif f["modified_time"] != prev.get("modified_time"):
            # Modified time changed but md5 same or unavailable (native Google files)
            if f["is_native_google"]:
                changes.append({**f, "change_type": "MODIFIED"})
            else:
                changes.append({**f, "change_type": "METADATA_CHANGED"})
        elif f["name"] != prev.get("name"):
            changes.append({
                **f,
                "change_type": "RENAMED",
                "previous_name": prev["name"],
            })
        else:
            changes.append({**f, "change_type": "UNCHANGED"})

    # Detect deletions
    for fid, prev in previous_files_by_id.items():
        if fid not in current_ids:
            changes.append({**prev, "change_type": "DELETED"})

    return changes


# ──────────────────────────────────────────────────────────
# Drive Activity API
# ──────────────────────────────────────────────────────────

def _parse_activity_response(resp):
    """Parse activity API response into structured records."""
    records = []
    for activity in resp.get("activities", []):
        timestamp = activity.get("timestamp") or ""
        if not timestamp:
            ts_range = activity.get("timeRange", {})
            timestamp = ts_range.get("endTime") or ts_range.get("startTime", "")

        actors = []
        for actor in activity.get("actors", []):
            user = actor.get("user", {})
            known_user = user.get("knownUser", {})
            actors.append({
                "person_name": known_user.get("personName", ""),
                "is_current_user": known_user.get("isCurrentUser", False),
            })

        actions = []
        for action in activity.get("actions", []):
            detail = action.get("detail", {})
            action_type = next(iter(detail.keys()), "unknown") if detail else "unknown"
            actions.append({
                "type": action_type,
                "detail": detail.get(action_type, {}),
            })

        targets = []
        for target in activity.get("targets", []):
            drive_item = target.get("driveItem", {})
            targets.append({
                "title": drive_item.get("title", ""),
                "name": drive_item.get("name", ""),
                "mime_type": drive_item.get("mimeType", ""),
                "file_id": drive_item.get("name", "").replace("items/", ""),
            })

        records.append({
            "timestamp": timestamp,
            "actors": actors,
            "actions": actions,
            "targets": targets,
        })
    return records


def fetch_recent_activity(activity_service, folder_id, since_timestamp=None,
                          file_ids=None):
    """Fetch recent activity from Drive Activity API.

    First tries folder-level query (ancestorName). If that fails (e.g. no ownership),
    falls back to per-file queries (itemName) for changed files.

    Args:
        since_timestamp: ISO8601 timestamp to filter events from. If None, fetches from 2020.
        file_ids: List of (file_id, file_name) tuples for per-file fallback.
    """
    time_filter = f"time >= '{since_timestamp}'" if since_timestamp else "time >= '2020-01-01T00:00:00Z'"

    # Strategy 1: Folder-level query (works if you own or have editor access)
    activities = _fetch_activity_by_ancestor(activity_service, folder_id, time_filter)
    if activities is not None:
        return activities

    # Strategy 2: Per-file queries (works with read-only access)
    if file_ids:
        print(f"    Falling back to per-file activity queries ({len(file_ids)} files)...")
        return _fetch_activity_by_files(activity_service, file_ids, time_filter)

    return []


def _fetch_activity_by_ancestor(activity_service, folder_id, time_filter):
    """Try folder-level activity query. Returns None if access denied."""
    activities = []
    try:
        body = {
            "ancestorName": f"items/{folder_id}",
            "pageSize": 100,
            "filter": time_filter,
        }
        page_token = None

        while True:
            if page_token:
                body["pageToken"] = page_token
            resp = activity_service.activity().query(body=body).execute()
            activities.extend(_parse_activity_response(resp))
            page_token = resp.get("nextPageToken")
            if not page_token:
                break

        return activities

    except Exception as e:
        print(f"    Folder-level activity query failed (likely no ownership): {e}")
        return None  # Signal to caller to try per-file fallback


def _fetch_activity_by_files(activity_service, file_ids, time_filter):
    """Query activity per-file using itemName. Works with read-only access."""
    all_activities = []
    seen_timestamps = set()  # Dedup across files

    for file_id, file_name in file_ids:
        try:
            body = {
                "itemName": f"items/{file_id}",
                "pageSize": 50,
                "filter": time_filter,
            }
            resp = activity_service.activity().query(body=body).execute()
            records = _parse_activity_response(resp)

            for record in records:
                # Dedup by timestamp + first target
                dedup_key = record["timestamp"] + str(record.get("targets", [{}])[0].get("file_id", ""))
                if dedup_key not in seen_timestamps:
                    seen_timestamps.add(dedup_key)
                    all_activities.append(record)

            time.sleep(0.05)  # Light rate limiting

        except Exception as e:
            # Skip individual file failures silently
            pass

    # Sort by timestamp descending
    all_activities.sort(key=lambda a: a.get("timestamp", ""), reverse=True)
    return all_activities


# ──────────────────────────────────────────────────────────
# Revision History
# ──────────────────────────────────────────────────────────

REVISION_HISTORY_FILE = BASE_DIR / "state" / "revision_history.json"
ACTIVITY_HISTORY_FILE = BASE_DIR / "state" / "activity_history.json"


def fetch_file_revisions(drive_service, file_id, file_name):
    """Fetch full revision history for a file.
    Returns list of {id, modifiedTime, lastModifyingUser, size}."""
    revisions = []
    try:
        page_token = None
        while True:
            resp = drive_service.revisions().list(
                fileId=file_id,
                fields="nextPageToken,revisions(id,modifiedTime,lastModifyingUser,size)",
                pageSize=1000,
                pageToken=page_token,
            ).execute()

            for rev in resp.get("revisions", []):
                user = rev.get("lastModifyingUser", {})
                revisions.append({
                    "id": rev.get("id", ""),
                    "time": rev.get("modifiedTime", ""),
                    "user_email": user.get("emailAddress", ""),
                    "user_name": user.get("displayName", ""),
                    "size": int(rev.get("size", 0) or 0),
                })

            page_token = resp.get("nextPageToken")
            if not page_token:
                break

    except Exception as e:
        print(f"    Warning: Could not fetch revisions for {file_name}: {e}")

    return revisions


def load_revision_history():
    """Load cumulative revision history from state file."""
    if REVISION_HISTORY_FILE.exists():
        try:
            with open(REVISION_HISTORY_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    return {"last_updated": "", "files": {}}


def save_revision_history(history):
    """Save cumulative revision history to state file."""
    REVISION_HISTORY_FILE.parent.mkdir(parents=True, exist_ok=True)
    history["last_updated"] = datetime.now(timezone.utc).isoformat()
    with open(REVISION_HISTORY_FILE, "w", encoding="utf-8") as f:
        json.dump(history, f, indent=2, ensure_ascii=False)


def load_activity_history():
    """Load cumulative activity history from state file."""
    if ACTIVITY_HISTORY_FILE.exists():
        try:
            with open(ACTIVITY_HISTORY_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    return {"last_updated": "", "events": []}


def save_activity_history(history):
    """Save cumulative activity history, capping at 1000 events."""
    ACTIVITY_HISTORY_FILE.parent.mkdir(parents=True, exist_ok=True)
    history["last_updated"] = datetime.now(timezone.utc).isoformat()
    # Cap at 1000 most recent events
    if len(history["events"]) > 1000:
        history["events"] = history["events"][-1000:]
    with open(ACTIVITY_HISTORY_FILE, "w", encoding="utf-8") as f:
        json.dump(history, f, indent=2, ensure_ascii=False)


# ──────────────────────────────────────────────────────────
# File downloading
# ──────────────────────────────────────────────────────────

EXPORT_MIMES = {
    "application/vnd.google-apps.document":
        ("application/vnd.openxmlformats-officedocument.wordprocessingml.document", ".docx"),
    "application/vnd.google-apps.spreadsheet":
        ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", ".xlsx"),
    "application/vnd.google-apps.presentation":
        ("application/vnd.openxmlformats-officedocument.presentationml.presentation", ".pptx"),
}


SKIP_MIMES = {
    "application/vnd.google-apps.shortcut",
    "application/vnd.google-apps.folder",
    "application/vnd.google-apps.form",
    "application/vnd.google-apps.map",
    "application/vnd.google-apps.site",
}


def _download_from_exports_folder(service, source_file_id, dest_dir, file_name):
    """Try to download a pre-exported PPTX from the Apps Script exports folder.
    Returns (path, md5) or (None, None) if not found."""
    if not EXPORTS_FOLDER_ID:
        return None, None

    try:
        # Search exports folder (and subfolders) for a file whose description
        # contains this source file ID
        query = (
            f"'{EXPORTS_FOLDER_ID}' in parents or "
            f"'{EXPORTS_FOLDER_ID}' in parents"
        )
        # Search recursively: look for PPTX files with source_id in description
        resp = service.files().list(
            q=(
                f"mimeType = 'application/vnd.openxmlformats-officedocument.presentationml.presentation' "
                f"and trashed = false "
                f"and fullText contains '{source_file_id}'"
            ),
            fields="files(id,name,description,size,md5Checksum)",
            pageSize=10,
        ).execute()

        for export_file in resp.get("files", []):
            desc = export_file.get("description", "")
            if source_file_id in desc:
                # Found the pre-exported PPTX — download it as binary
                print(f"    Found pre-exported PPTX in exports folder: {export_file['name']}")
                export_id = export_file["id"]
                request = service.files().get_media(fileId=export_id)

                dest_dir = Path(dest_dir)
                dest_dir.mkdir(parents=True, exist_ok=True)
                dest_name = file_name if file_name.endswith(".pptx") else file_name + ".pptx"
                dest_path = dest_dir / dest_name

                fh = io.BytesIO()
                downloader = MediaIoBaseDownload(fh, request)
                done = False
                while not done:
                    _, done = downloader.next_chunk()

                with open(dest_path, "wb") as f:
                    f.write(fh.getvalue())

                local_md5 = hashlib.md5(fh.getvalue()).hexdigest()
                return str(dest_path), local_md5

    except Exception as e:
        print(f"    Exports folder lookup failed: {e}")

    return None, None


def download_file(service, file_meta, dest_dir):
    """Download a file from Drive. Exports native Google formats to Office equivalents.
    For large native Slides that exceed the 10MB export limit, falls back to
    pre-exported PPTX files from the Apps Script exports folder."""
    fid = file_meta["id"]
    name = file_meta["name"]
    mime = file_meta["mime_type"]

    # Skip non-downloadable types (shortcuts, folders, forms, etc.)
    if mime in SKIP_MIMES:
        print(f"    Skipping {name} (type: {mime.split('.')[-1]})")
        return None, None

    dest_dir = Path(dest_dir)
    dest_dir.mkdir(parents=True, exist_ok=True)

    # For native Google files, export to Office format
    if mime in EXPORT_MIMES:
        export_mime, ext = EXPORT_MIMES[mime]
        if not name.endswith(ext):
            name = name + ext
        request = service.files().export_media(fileId=fid, mimeType=export_mime)
    else:
        request = service.files().get_media(fileId=fid)

    dest_path = dest_dir / name
    fh = io.BytesIO()
    downloader = MediaIoBaseDownload(fh, request)

    try:
        done = False
        while not done:
            _, done = downloader.next_chunk()
    except Exception as e:
        if "exportSizeLimitExceeded" in str(e):
            # Try the Apps Script exports folder fallback
            print(f"    Export too large ({name}), checking exports folder...")
            fallback_path, fallback_md5 = _download_from_exports_folder(
                service, fid, dest_dir, name
            )
            if fallback_path:
                return fallback_path, fallback_md5
            # No fallback available — re-raise so caller logs the error
            print(f"    No pre-exported PPTX found. Run Apps Script to export.")
        raise

    with open(dest_path, "wb") as f:
        f.write(fh.getvalue())

    # Compute local MD5
    local_md5 = hashlib.md5(fh.getvalue()).hexdigest()

    return str(dest_path), local_md5


# ──────────────────────────────────────────────────────────
# PPTX Integrity Verification
# ──────────────────────────────────────────────────────────

def verify_downloaded_pptx(sources_dir):
    """Verify all downloaded PPTX files can be opened and have valid content.
    Returns {valid, total, errors: [{file, error}], warnings: [{file, warning}]}."""
    results = {"valid": 0, "total": 0, "errors": [], "warnings": []}

    try:
        from pptx import Presentation
    except ImportError:
        print("  [WARN] python-pptx not installed, skipping PPTX integrity check")
        return results

    pptx_files = list(Path(sources_dir).rglob("*.pptx"))
    results["total"] = len(pptx_files)

    if not pptx_files:
        return results

    print(f"\n  Verifying {len(pptx_files)} PPTX files...")

    for pptx_path in pptx_files:
        rel_path = str(pptx_path.relative_to(sources_dir))
        try:
            prs = Presentation(str(pptx_path))
            slide_count = len(prs.slides)
            total_shapes = 0
            total_images = 0

            for slide in prs.slides:
                for shape in slide.shapes:
                    total_shapes += 1
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        total_images += 1

            # Warnings for suspicious files
            if slide_count == 0:
                results["warnings"].append({
                    "file": rel_path,
                    "warning": "Empty presentation (0 slides)",
                })
            elif total_shapes == 0:
                results["warnings"].append({
                    "file": rel_path,
                    "warning": f"{slide_count} slides but 0 shapes (possibly corrupt)",
                })
            else:
                results["valid"] += 1

        except Exception as e:
            results["errors"].append({
                "file": rel_path,
                "error": str(e)[:200],
            })

    print(f"  PPTX integrity: {results['valid']}/{results['total']} valid, "
          f"{len(results['errors'])} corrupt, {len(results['warnings'])} warnings")

    return results


# ──────────────────────────────────────────────────────────
# Main sync
# ──────────────────────────────────────────────────────────

def run_sync(dry_run=False, download_all=False):
    """Execute full sync: scan → detect changes → download → log.

    Args:
        dry_run: If True, scan and report changes but skip downloads
                 and do NOT update previous_scan.json.
        download_all: If True, download ALL files (not just changed ones).
                      Needed in CI where the runner starts with an empty sources dir.
    """
    print("=" * 60)
    print("  KB Maintenance Pipeline — Drive Sync" + (" (DRY RUN)" if dry_run else ""))
    print("=" * 60)
    print()

    # Ensure directories
    SOURCES_DIR.mkdir(parents=True, exist_ok=True)
    LOGS_DIR.mkdir(parents=True, exist_ok=True)
    PREVIOUS_SCAN_FILE.parent.mkdir(parents=True, exist_ok=True)

    # Authenticate
    creds = authenticate()
    drive_service = get_drive_service(creds)
    activity_service = get_activity_service(creds)
    print("Authenticated.\n")

    # Load previous scan
    prev_scan = load_previous_scan()
    prev_files_by_id = {}
    for term_data in prev_scan.get("terms", {}).values():
        for f in term_data.get("files", []):
            prev_files_by_id[f["id"]] = f

    # Determine last sync timestamp for incremental activity queries
    last_sync_timestamp = prev_scan.get("scan_timestamp")

    # Scan all target folders
    sync_result = {
        "sync_timestamp": datetime.now(timezone.utc).isoformat(),
        "dry_run": dry_run,
        "terms": {},
        "summary": {
            "total_files": 0,
            "new": 0,
            "modified": 0,
            "deleted": 0,
            "renamed": 0,
            "metadata_changed": 0,
            "unchanged": 0,
            "downloaded": 0,
            "errors": 0,
        },
        "activity_log": {},
        "revision_history": {},
    }

    all_current_files = {}  # For saving as next "previous scan"

    # Load cumulative histories
    revision_history = load_revision_history()
    activity_history = load_activity_history()

    for term_key, folder_info in TARGET_FOLDERS.items():
        folder_id = folder_info["id"]
        folder_name = folder_info["name"]
        print(f"--- Scanning: {folder_name} ({term_key}) ---")

        # Scan
        current_files = scan_folder(drive_service, folder_id)
        print(f"  Found {len(current_files)} files")

        # Detect changes
        term_prev = {
            f["id"]: f
            for f in prev_scan.get("terms", {}).get(term_key, {}).get("files", [])
        }
        changes = detect_changes(current_files, term_prev)

        # Fetch activity (incremental: since last sync if available)
        print(f"  Fetching activity log...")
        # For per-file fallback, only query changed files (not all files)
        changed_ids = [
            (c["id"], c["name"]) for c in changes
            if c["change_type"] in ("NEW", "MODIFIED", "METADATA_CHANGED", "RENAMED")
        ]
        activities = fetch_recent_activity(
            activity_service, folder_id,
            since_timestamp=last_sync_timestamp,
            file_ids=changed_ids if changed_ids else None,
        )
        sync_result["activity_log"][term_key] = activities
        print(f"  {len(activities)} activity records")

        # Append to cumulative activity history
        for act in activities:
            actors_str = ", ".join(
                a.get("person_name", "unknown") for a in act.get("actors", [])
            )
            for action in act.get("actions", []):
                for target in act.get("targets", []):
                    activity_history["events"].append({
                        "timestamp": act.get("timestamp", ""),
                        "actor": actors_str,
                        "action": action.get("type", "unknown"),
                        "target": target.get("title", ""),
                        "term": term_key,
                    })

        # Fetch revision history for changed files
        changed_files = [
            c for c in changes
            if c["change_type"] in ("NEW", "MODIFIED", "METADATA_CHANGED")
        ]
        if changed_files:
            print(f"  Fetching revision history for {len(changed_files)} changed files...")
            for cf in changed_files:
                revisions = fetch_file_revisions(drive_service, cf["id"], cf["name"])
                if revisions:
                    file_id = cf["id"]
                    sync_result["revision_history"][file_id] = {
                        "name": cf["name"],
                        "term": term_key,
                        "revisions": revisions,
                    }
                    # Update cumulative revision history
                    revision_history["files"][file_id] = {
                        "name": cf["name"],
                        "term": term_key,
                        "revisions": revisions,
                    }
                time.sleep(0.1)  # Rate limiting between revision queries

        # Download changed files (skip in dry-run mode)
        term_sources = SOURCES_DIR / term_key
        downloaded = []
        errors = []

        for change in changes:
            ct = change["change_type"]
            sync_result["summary"]["total_files"] += 1

            if ct == "UNCHANGED":
                sync_result["summary"]["unchanged"] += 1
            elif ct == "NEW":
                sync_result["summary"]["new"] += 1
            elif ct == "MODIFIED":
                sync_result["summary"]["modified"] += 1
            elif ct == "DELETED":
                sync_result["summary"]["deleted"] += 1
            elif ct == "RENAMED":
                sync_result["summary"]["renamed"] += 1
            elif ct == "METADATA_CHANGED":
                sync_result["summary"]["metadata_changed"] += 1

            # Determine if this file should be downloaded
            should_download = (
                not dry_run and (
                    ct in ("NEW", "MODIFIED", "RENAMED") or
                    (download_all and ct in ("UNCHANGED", "METADATA_CHANGED"))
                )
            )

            if should_download:
                try:
                    # Preserve folder hierarchy from Drive
                    fp = change.get("folder_path", "")
                    dest_dir = term_sources / fp if fp else term_sources
                    label = ct if ct != "UNCHANGED" else "SYNC"
                    print(f"  [{label}] Downloading: {fp}/{change['name']}" if fp else f"  [{label}] Downloading: {change['name']}")
                    local_path, local_md5 = download_file(
                        drive_service, change, dest_dir
                    )
                    if local_path is None:
                        # File was skipped (shortcut, folder, etc.)
                        continue
                    change["local_path"] = local_path
                    change["local_md5"] = local_md5
                    downloaded.append(change["name"])
                    sync_result["summary"]["downloaded"] += 1
                except Exception as e:
                    print(f"    ERROR downloading {change['name']}: {e}")
                    change["download_error"] = str(e)
                    errors.append({"file": change["name"], "error": str(e)})
                    sync_result["summary"]["errors"] += 1
            elif ct in ("NEW", "MODIFIED", "RENAMED") and dry_run:
                print(f"  [{ct}] Would download: {change['name']}")

        # Store term results
        sync_result["terms"][term_key] = {
            "folder_id": folder_id,
            "folder_name": folder_name,
            "files": changes,
            "downloaded": downloaded,
            "errors": errors,
        }

        # Save for next comparison
        all_current_files[term_key] = {
            "files": current_files,
        }

        if dry_run:
            print(f"  Would download: {sum(1 for c in changes if c['change_type'] in ('NEW', 'MODIFIED', 'RENAMED'))} files")
        else:
            print(f"  Downloaded: {len(downloaded)} files")
        if errors:
            print(f"  Errors: {len(errors)}")
        print()

    # Verify integrity of downloaded PPTX files
    if not dry_run:
        integrity_results = verify_downloaded_pptx(SOURCES_DIR)
        sync_result["integrity"] = integrity_results
        if integrity_results["errors"] or integrity_results["warnings"]:
            print(f"\n  PPTX Integrity: {integrity_results['valid']} valid, "
                  f"{len(integrity_results['errors'])} errors, "
                  f"{len(integrity_results['warnings'])} warnings")

    # Save current scan as "previous" for next run (skip in dry-run mode)
    if not dry_run:
        with open(PREVIOUS_SCAN_FILE, "w", encoding="utf-8") as f:
            json.dump(
                {"scan_timestamp": sync_result["sync_timestamp"], "terms": all_current_files},
                f, indent=2, ensure_ascii=False,
            )

    # Save cumulative histories (always, even in dry-run for audit trail)
    save_revision_history(revision_history)
    save_activity_history(activity_history)

    # Write comprehensive sync log (append-style: one entry per sync)
    log_entry = sync_result
    log_file = LOGS_DIR / f"sync_{datetime.now(timezone.utc).strftime('%Y%m%d_%H%M%S')}.json"
    with open(log_file, "w", encoding="utf-8") as f:
        json.dump(log_entry, f, indent=2, ensure_ascii=False)

    # Summary
    s = sync_result["summary"]
    print("=" * 60)
    print("  Sync Complete" + (" (DRY RUN)" if dry_run else ""))
    print("=" * 60)
    print(f"  Total files scanned : {s['total_files']}")
    print(f"  New                 : {s['new']}")
    print(f"  Modified            : {s['modified']}")
    print(f"  Deleted             : {s['deleted']}")
    print(f"  Renamed             : {s['renamed']}")
    print(f"  Metadata changed    : {s['metadata_changed']}")
    print(f"  Unchanged           : {s['unchanged']}")
    if not dry_run:
        print(f"  Downloaded          : {s['downloaded']}")
    print(f"  Errors              : {s['errors']}")
    if dry_run:
        print(f"\n  ** DRY RUN — no files downloaded, previous_scan.json NOT updated **")
    print(f"\n  Log saved: {log_file}")
    print("=" * 60)

    return sync_result


if __name__ == "__main__":
    run_sync()
