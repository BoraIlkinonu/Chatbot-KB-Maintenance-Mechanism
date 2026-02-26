"""
QA Runner — Orchestrates all 4 layers and produces unified report.
"""

import sys
import json
from pathlib import Path

sys.stdout.reconfigure(encoding="utf-8")

from config import OUTPUT_DIR, CONSOLIDATED_DIR, VALIDATION_DIR, MEDIA_DIR, SOURCES_DIR
from qa.report import QAReport, CheckResult
from qa.config import TERM_PROFILES


def _load_kb(term: int) -> dict | None:
    path = OUTPUT_DIR / f"Term {term} - Lesson Based Structure.json"
    if path.exists():
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)
    return None


def _find_source_pptx_for_lesson(term: int, lesson_num: int) -> Path | None:
    """Find source PPTX for a lesson (for Layer 2)."""
    cons_path = CONSOLIDATED_DIR / f"consolidated_term{term}.json"
    if not cons_path.exists():
        return None
    with open(cons_path, "r", encoding="utf-8") as f:
        cons = json.load(f)
    lesson_data = cons.get("by_lesson", {}).get(str(lesson_num), {})
    for doc in lesson_data.get("documents", []):
        rel_path = doc.get("path", "")
        source = SOURCES_DIR / Path(rel_path).with_suffix(".pptx")
        if source.exists():
            return source
    return None


def _extract_pptx_ground_truth(pptx_path: Path) -> dict:
    """Extract ground truth from PPTX (imported from cross_validate_kb)."""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slides = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_data = {"slide_number": slide_num, "text": [], "notes": "", "links": [], "tables": []}
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_data["text"].append(shape.text.strip())
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                if run.hyperlink and run.hyperlink.address:
                                    slide_data["links"].append({"url": run.hyperlink.address, "text": run.text.strip()})
                            except (KeyError, TypeError, AttributeError):
                                pass
                if shape.has_table:
                    table_data = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                    if table_data:
                        slide_data["tables"].append(table_data)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_data["notes"] = notes
            slides.append(slide_data)
        return {"file": str(pptx_path), "total_slides": len(slides), "slides": slides}
    except Exception as e:
        return {"error": str(e), "slides": []}


def run_qa(layers=None, terms=None, budget=None, verbose=False, json_output=False) -> QAReport:
    """Run the QA system.

    Args:
        layers: List of layer numbers to run (default: all)
        terms: List of term numbers to check (default: all)
        budget: LLM call budget for Layer 2
        verbose: Print detailed per-check output
        json_output: Suppress text output, return JSON-compatible report

    Returns:
        QAReport with all results
    """
    if layers is None:
        layers = [1, 2, 3, 4]
    if terms is None:
        terms = [1, 2, 3]

    report = QAReport()
    previous_builds_dir = VALIDATION_DIR / "previous_builds"

    if not json_output:
        print("=" * 60)
        print("  QA System — Layered Validation")
        print("=" * 60)
        print()

    # ── Layer 1: Programmatic Validation ──
    if 1 in layers:
        if not json_output:
            print(">>> Layer 1: Programmatic Validation")
            print("-" * 40)

        from qa.layer1.schema_checks import run_schema_checks
        from qa.layer1.content_checks import run_content_checks
        from qa.layer1.consistency_checks import run_consistency_checks
        from qa.layer1.regression_checks import run_regression_checks

        for term in terms:
            kb = _load_kb(term)
            if not kb:
                report.add(CheckResult(
                    check_id="L1_LOAD", layer=1, severity="ERROR",
                    passed=False, message=f"Cannot load KB for term {term}",
                    details={"term": term},
                ))
                continue

            if not json_output:
                print(f"  Term {term}:")

            # Schema
            schema_results = run_schema_checks(kb, term)
            report.add_all(schema_results)

            # Content
            content_results = run_content_checks(kb, term)
            report.add_all(content_results)

            # Consistency
            consistency_results = run_consistency_checks(kb, term, CONSOLIDATED_DIR, OUTPUT_DIR)
            report.add_all(consistency_results)

            # Regression
            regression_results = run_regression_checks(kb, term, previous_builds_dir)
            report.add_all(regression_results)

            l1_failures = sum(1 for r in schema_results + content_results + consistency_results + regression_results if not r.passed)
            if not json_output:
                total = len(schema_results) + len(content_results) + len(consistency_results) + len(regression_results)
                print(f"    {total} checks, {l1_failures} failures")

        # Cross-term check X007 (deferred from per-term)
        all_kbs = {t: _load_kb(t) for t in terms}
        cross_term_issues = []
        for lid in range(1, 23):
            titles = {}
            for t, kb in all_kbs.items():
                if not kb:
                    continue
                for l in kb.get("lessons", []):
                    if l.get("metadata", {}).get("lesson_id") == lid:
                        titles[t] = l.get("lesson_title", "").strip().lower()
            # Check if same lesson_id has identical titles across different terms
            title_values = list(titles.values())
            if len(title_values) > 1 and len(set(title_values)) < len(title_values):
                cross_term_issues.append({"lesson_id": lid, "titles": titles})
        report.add(CheckResult(
            check_id="X007", layer=1, severity="WARNING",
            passed=len(cross_term_issues) == 0,
            message=f"{len(cross_term_issues)} identical titles across terms for same lesson_id" if cross_term_issues else "No cross-term title collisions",
            details={"issues": cross_term_issues},
        ))
        # Update the deferred X007 results
        report.results = [r for r in report.results if not (r.check_id == "X007" and r.details.get("deferred"))]

        if not json_output:
            l1_total = len(report.by_layer(1))
            l1_fail = sum(1 for r in report.by_layer(1) if not r.passed)
            print(f"  Layer 1 total: {l1_total} checks, {l1_fail} failures")
            print()

    # ── Layer 2: LLM Cross-Validation ──
    if 2 in layers:
        if not json_output:
            print(">>> Layer 2: LLM Cross-Validation")
            print("-" * 40)

        # Prepare data for Layer 2
        l1_errors = [r for r in report.by_layer(1) if not r.passed and r.severity in ("ERROR", "WARNING")]

        # Build lessons_with_sources for Phase 2
        lessons_with_sources = []
        for term in terms:
            kb = _load_kb(term)
            if not kb:
                continue
            for l in kb.get("lessons", []):
                lid = l.get("metadata", {}).get("lesson_id", 0)
                pptx_path = _find_source_pptx_for_lesson(term, lid)
                if pptx_path:
                    source_content = _extract_pptx_ground_truth(pptx_path)
                    lessons_with_sources.append({
                        "term": term,
                        "lesson_num": lid,
                        "source_content": source_content,
                        "kb_lesson": l,
                    })

        from qa.layer2.investigator import run_layer2
        l2_results, l2_summary = run_layer2(
            errors=l1_errors,
            lessons_with_sources=lessons_with_sources,
            terms=terms,
            budget=budget,
            verbose=verbose,
        )
        report.add_all(l2_results)
        report.set_layer_summary(2, l2_summary)
        report.set_llm_meta({
            "model": l2_summary.get("model", "sonnet"),
            "budget": l2_summary.get("budget"),
            "calls_made": l2_summary.get("calls_made"),
        })

        if not json_output:
            print(f"  Layer 2: {len(l2_results)} checks, confidence: {l2_summary.get('confidence', 'N/A')}")
            print()

    # ── Layer 3: Edge Case Tests ──
    if 3 in layers:
        if not json_output:
            print(">>> Layer 3: Edge Case Tests")
            print("-" * 40)

        from qa.layer3.real_output_tests import run_real_output_tests
        from qa.layer3.boundary_tests import run_boundary_tests
        from qa.layer3.integrity_tests import run_integrity_tests

        l3_results = []
        l3_results.extend(run_real_output_tests(OUTPUT_DIR, MEDIA_DIR))
        l3_results.extend(run_boundary_tests(OUTPUT_DIR))
        l3_results.extend(run_integrity_tests(OUTPUT_DIR, CONSOLIDATED_DIR))
        report.add_all(l3_results)

        if not json_output:
            l3_fail = sum(1 for r in l3_results if not r.passed)
            print(f"  Layer 3: {len(l3_results)} checks, {l3_fail} failures")
            print()

    # ── Layer 4: User Scenario Tests ──
    if 4 in layers:
        if not json_output:
            print(">>> Layer 4: User Scenario Tests")
            print("-" * 40)

        from qa.layer4.retrieval_tests import run_retrieval_tests
        from qa.layer4.completeness_tests import run_completeness_tests
        from qa.layer4.navigability_tests import run_navigability_tests

        l4_results = []
        l4_results.extend(run_retrieval_tests(OUTPUT_DIR))
        l4_results.extend(run_completeness_tests(OUTPUT_DIR))
        l4_results.extend(run_navigability_tests(OUTPUT_DIR))
        report.add_all(l4_results)

        if not json_output:
            l4_fail = sum(1 for r in l4_results if not r.passed)
            print(f"  Layer 4: {len(l4_results)} checks, {l4_fail} failures")
            print()

    # ── Save Reports ──
    VALIDATION_DIR.mkdir(parents=True, exist_ok=True)

    report.save(VALIDATION_DIR / "qa_report.json")
    report.save_text(VALIDATION_DIR / "qa_report.txt")

    # Per-term reports
    for term in terms:
        term_results = [r for r in report.results if r.details.get("term") == term]
        if term_results:
            term_report = QAReport()
            term_report.results = term_results
            term_report.layer_summaries = report.layer_summaries
            term_report.save(VALIDATION_DIR / f"qa_report_term{term}.json")

    # Archive current builds for regression
    from qa.layer1.regression_checks import archive_current_build
    verdict = report.compute_verdict()
    if verdict == "PASS":
        archive_current_build(OUTPUT_DIR, previous_builds_dir)
        if not json_output:
            print("  Archived current builds for future regression checks")

    if not json_output:
        print()
        report.print_summary(verbose)

    return report
